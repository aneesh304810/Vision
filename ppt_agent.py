"""
╔══════════════════════════════════════════════════════════════════════╗
║        AI PRESENTATION INTELLIGENCE PLATFORM                        ║
║        Powered by Qwen2.5 via llama.cpp (OpenAI-compatible API)     ║
╚══════════════════════════════════════════════════════════════════════╝

Run:
    pip install streamlit python-pptx requests pillow
    streamlit run ppt_agent.py

Requirements:
    llama.cpp server running locally with Qwen2.5-Coder-14B
    Default endpoint: http://localhost:8080/v1/chat/completions
"""

import io
import json
import base64
import zipfile
import textwrap
import requests
import tempfile
import traceback
from pathlib import Path
from copy import deepcopy

import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ─── PAGE CONFIG ─────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="PPT Intelligence Platform",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─── CUSTOM CSS ──────────────────────────────────────────────────────────────

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Serif+Display&family=DM+Mono:wght@400;500&family=DM+Sans:wght@300;400;500;600&display=swap');

/* Root */
:root {
    --bg: #0d0f14;
    --surface: #161920;
    --surface2: #1e2230;
    --border: #2a2f42;
    --accent: #4f8ef7;
    --accent2: #7c3aed;
    --success: #22c55e;
    --warning: #f59e0b;
    --danger: #ef4444;
    --text: #e8eaf0;
    --muted: #6b7280;
    --font-display: 'DM Serif Display', serif;
    --font-body: 'DM Sans', sans-serif;
    --font-mono: 'DM Mono', monospace;
}

html, body, [class*="css"] {
    font-family: var(--font-body) !important;
    background-color: var(--bg) !important;
    color: var(--text) !important;
}

/* Hide Streamlit chrome */
#MainMenu, footer, header { visibility: hidden; }
.block-container { padding: 2rem 2rem 4rem !important; max-width: 1400px !important; }

/* Sidebar */
[data-testid="stSidebar"] {
    background: var(--surface) !important;
    border-right: 1px solid var(--border) !important;
}
[data-testid="stSidebar"] * { color: var(--text) !important; }

/* Hero header */
.hero {
    background: linear-gradient(135deg, #0d0f14 0%, #161b2e 50%, #0d0f14 100%);
    border: 1px solid var(--border);
    border-radius: 16px;
    padding: 2.5rem 3rem;
    margin-bottom: 2rem;
    position: relative;
    overflow: hidden;
}
.hero::before {
    content: '';
    position: absolute;
    top: -50%;
    left: -20%;
    width: 60%;
    height: 200%;
    background: radial-gradient(ellipse, rgba(79,142,247,0.08) 0%, transparent 70%);
    pointer-events: none;
}
.hero-title {
    font-family: var(--font-display) !important;
    font-size: 2.4rem !important;
    font-weight: 400 !important;
    color: #fff !important;
    margin: 0 0 0.4rem !important;
    line-height: 1.1 !important;
}
.hero-sub {
    font-size: 0.95rem;
    color: var(--muted);
    font-weight: 300;
    letter-spacing: 0.02em;
}
.hero-badge {
    display: inline-block;
    background: rgba(79,142,247,0.12);
    border: 1px solid rgba(79,142,247,0.3);
    color: var(--accent);
    padding: 0.2rem 0.75rem;
    border-radius: 100px;
    font-size: 0.75rem;
    font-family: var(--font-mono);
    margin-bottom: 1rem;
    letter-spacing: 0.05em;
}

/* Cards */
.card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 12px;
    padding: 1.5rem;
    margin-bottom: 1rem;
}
.card-header {
    font-size: 0.7rem;
    font-family: var(--font-mono);
    color: var(--muted);
    letter-spacing: 0.1em;
    text-transform: uppercase;
    margin-bottom: 0.75rem;
}

/* Slide cards */
.slide-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 12px;
    padding: 1.5rem 1.75rem;
    margin-bottom: 1rem;
    transition: border-color 0.2s;
    position: relative;
    overflow: hidden;
}
.slide-card::before {
    content: '';
    position: absolute;
    left: 0; top: 0; bottom: 0;
    width: 3px;
    background: var(--accent);
    border-radius: 3px 0 0 3px;
}
.slide-card:hover { border-color: rgba(79,142,247,0.4); }

.slide-num {
    font-family: var(--font-mono);
    font-size: 0.7rem;
    color: var(--accent);
    letter-spacing: 0.1em;
    margin-bottom: 0.4rem;
}
.slide-title {
    font-family: var(--font-display) !important;
    font-size: 1.15rem !important;
    color: #fff !important;
    margin-bottom: 0.6rem !important;
    line-height: 1.3 !important;
}
.slide-meta {
    display: flex;
    gap: 0.5rem;
    flex-wrap: wrap;
    margin-bottom: 0.75rem;
}
.tag {
    font-size: 0.68rem;
    font-family: var(--font-mono);
    padding: 0.15rem 0.6rem;
    border-radius: 100px;
    letter-spacing: 0.05em;
}
.tag-layout { background: rgba(124,58,237,0.15); border: 1px solid rgba(124,58,237,0.3); color: #a78bfa; }
.tag-class  { background: rgba(79,142,247,0.12); border: 1px solid rgba(79,142,247,0.3); color: var(--accent); }
.tag-visual { background: rgba(34,197,94,0.1);  border: 1px solid rgba(34,197,94,0.25); color: #4ade80; }

.bullet-list { margin: 0; padding: 0; list-style: none; }
.bullet-list li {
    font-size: 0.875rem;
    color: #c8cdd8;
    padding: 0.25rem 0;
    padding-left: 1.1rem;
    position: relative;
    line-height: 1.5;
}
.bullet-list li::before {
    content: '▸';
    position: absolute;
    left: 0;
    color: var(--accent);
    font-size: 0.65rem;
    top: 0.35rem;
}

/* Flow steps */
.flow-steps {
    display: flex;
    flex-wrap: wrap;
    align-items: center;
    gap: 0.4rem;
    margin-top: 0.75rem;
}
.flow-step {
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 6px;
    padding: 0.3rem 0.75rem;
    font-size: 0.78rem;
    font-family: var(--font-mono);
    color: var(--text);
}
.flow-arrow { color: var(--accent); font-size: 0.9rem; }

/* Metrics */
.metric-grid { display: flex; gap: 1rem; flex-wrap: wrap; margin-top: 0.75rem; }
.metric-card {
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 0.75rem 1.25rem;
    text-align: center;
    min-width: 90px;
}
.metric-value {
    font-family: var(--font-display);
    font-size: 1.6rem;
    color: var(--accent);
    line-height: 1;
}
.metric-label { font-size: 0.68rem; color: var(--muted); margin-top: 0.2rem; }

/* Two-column */
.two-col { display: grid; grid-template-columns: 1fr 1fr; gap: 1rem; margin-top: 0.75rem; }
.col-box {
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 0.75rem 1rem;
}
.col-heading { font-size: 0.78rem; font-weight: 600; color: #fff; margin-bottom: 0.4rem; }

/* Status bar */
.status-bar {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 0.6rem 1rem;
    font-family: var(--font-mono);
    font-size: 0.75rem;
    color: var(--muted);
    margin-bottom: 1rem;
    display: flex;
    gap: 1.5rem;
}
.status-item { display: flex; align-items: center; gap: 0.4rem; }
.status-dot { width: 6px; height: 6px; border-radius: 50%; background: var(--success); }
.status-dot.warn { background: var(--warning); }
.status-dot.off  { background: var(--muted); }

/* Log box */
.log-box {
    background: #080a0f;
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 1rem;
    font-family: var(--font-mono);
    font-size: 0.72rem;
    color: #6b7280;
    max-height: 200px;
    overflow-y: auto;
    line-height: 1.7;
}
.log-ok   { color: #22c55e; }
.log-info { color: #4f8ef7; }
.log-warn { color: #f59e0b; }
.log-err  { color: #ef4444; }

/* Buttons */
.stButton > button {
    background: var(--accent) !important;
    color: #fff !important;
    border: none !important;
    border-radius: 8px !important;
    font-family: var(--font-body) !important;
    font-weight: 500 !important;
    padding: 0.55rem 1.5rem !important;
    transition: opacity 0.2s !important;
}
.stButton > button:hover { opacity: 0.85 !important; }

/* Inputs */
.stTextInput > div > div > input,
.stNumberInput > div > div > input {
    background: var(--surface2) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    color: var(--text) !important;
    font-family: var(--font-mono) !important;
    font-size: 0.85rem !important;
}

/* Expander */
.streamlit-expanderHeader {
    background: var(--surface) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    color: var(--text) !important;
    font-family: var(--font-body) !important;
}

/* Progress */
.stProgress > div > div > div { background: var(--accent) !important; border-radius: 4px !important; }

/* Divider */
hr { border-color: var(--border) !important; margin: 1.5rem 0 !important; }

/* Deck header */
.deck-header {
    background: linear-gradient(90deg, rgba(79,142,247,0.1), transparent);
    border: 1px solid rgba(79,142,247,0.2);
    border-left: 3px solid var(--accent);
    border-radius: 8px;
    padding: 1rem 1.5rem;
    margin-bottom: 1.5rem;
}
.deck-title { font-family: var(--font-display); font-size: 1.4rem; color: #fff; }
.deck-meta { font-size: 0.78rem; color: var(--muted); margin-top: 0.25rem; font-family: var(--font-mono); }

/* Section label */
.section-label {
    font-size: 0.65rem;
    font-family: var(--font-mono);
    letter-spacing: 0.15em;
    color: var(--muted);
    text-transform: uppercase;
    margin: 1.5rem 0 0.75rem;
    display: flex;
    align-items: center;
    gap: 0.5rem;
}
.section-label::after {
    content: '';
    flex: 1;
    height: 1px;
    background: var(--border);
}
</style>
""", unsafe_allow_html=True)

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def log(msg: str, kind: str = "info"):
    cls = {"ok": "log-ok", "info": "log-info", "warn": "log-warn", "err": "log-err"}.get(kind, "log-info")
    prefix = {"ok": "✓", "info": "›", "warn": "⚠", "err": "✗"}.get(kind, "›")
    st.session_state.logs.append(f'<span class="{cls}">{prefix} {msg}</span>')

def render_logs():
    if st.session_state.logs:
        html = "<br>".join(st.session_state.logs[-30:])
        st.markdown(f'<div class="log-box">{html}</div>', unsafe_allow_html=True)

# ─── PPTX PARSER ─────────────────────────────────────────────────────────────

def parse_pptx(file_bytes: bytes) -> dict:
    """Extract structured content from a .pptx file."""
    prs = Presentation(io.BytesIO(file_bytes))
    slides_data = []
    images = {}

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": slide_idx + 1,
            "title": "",
            "bullets": [],
            "tables": [],
            "image_ids": [],
            "raw_text_blocks": [],
        }

        for shape in slide.shapes:
            # Title
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape.shape_type == 13:  # Picture
                    pass
                elif hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_idx = shape.placeholder_format.idx
                    if ph_idx == 0:  # Title
                        slide_info["title"] = text
                    elif ph_idx == 1:  # Body
                        for para in shape.text_frame.paragraphs:
                            para_text = para.text.strip()
                            if para_text:
                                slide_info["bullets"].append(para_text)
                    else:
                        slide_info["raw_text_blocks"].append(text)
                else:
                    # Non-placeholder text box
                    lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                    if lines:
                        if not slide_info["title"] and len(lines[0]) < 120:
                            slide_info["title"] = lines[0]
                            slide_info["raw_text_blocks"].extend(lines[1:])
                        else:
                            slide_info["raw_text_blocks"].extend(lines)

            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    slide_info["tables"].append(table_data)

            # Images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_id = f"img_{slide_idx}_{shape.shape_id}"
                    img_bytes = shape.image.blob
                    img_b64 = base64.b64encode(img_bytes).decode()
                    img_ext = shape.image.ext
                    images[img_id] = {
                        "data": img_b64,
                        "ext": img_ext,
                        "slide": slide_idx + 1,
                        "width": shape.width,
                        "height": shape.height,
                    }
                    slide_info["image_ids"].append(img_id)
                except Exception:
                    pass

        # Fallback title
        if not slide_info["title"] and slide_info["raw_text_blocks"]:
            slide_info["title"] = slide_info["raw_text_blocks"].pop(0)

        slides_data.append(slide_info)

    return {
        "slide_count": len(slides_data),
        "slides": slides_data,
        "images": images,
    }

# ─── DESIGN SYSTEM EXTRACTOR (single file) ───────────────────────────────────

def extract_design_system(file_bytes: bytes) -> dict:
    """Extract color palette, fonts, layout hints from a single PPTX."""
    prs = Presentation(io.BytesIO(file_bytes))
    colors, fonts = set(), set()
    bg_colors = set()
    title_fonts, body_fonts = set(), set()
    font_sizes = []

    for slide in prs.slides:
        # Background fill color
        try:
            bg = slide.background.fill
            if bg.type and bg.fore_color:
                bg_colors.add(str(bg.fore_color.rgb))
        except Exception:
            pass

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                            # Heuristic: large font → display/title font
                            if run.font.size and run.font.size >= Pt(24):
                                title_fonts.add(run.font.name)
                            elif run.font.size and run.font.size <= Pt(18):
                                body_fonts.add(run.font.name)
                        if run.font.size:
                            font_sizes.append(round(run.font.size / 12700))  # EMU → pt
                        if run.font.color and run.font.color.type:
                            try:
                                colors.add(str(run.font.color.rgb))
                            except Exception:
                                pass

            # Shape fill colors (background blocks, accent shapes)
            try:
                fill = shape.fill
                if fill.type == 1:  # solid
                    colors.add(str(fill.fore_color.rgb))
            except Exception:
                pass

    width_emu  = prs.slide_width
    height_emu = prs.slide_height
    aspect     = "16:9" if abs(width_emu / height_emu - 16 / 9) < 0.1 else "4:3"

    avg_font_size = round(sum(font_sizes) / len(font_sizes)) if font_sizes else 0
    max_font_size = max(font_sizes) if font_sizes else 0

    return {
        "aspect_ratio": aspect,
        "slide_width_inches":  round(width_emu  / 914400, 2),
        "slide_height_inches": round(height_emu / 914400, 2),
        "slide_count": len(prs.slides),
        "all_fonts":     sorted(fonts)[:12],
        "title_fonts":   sorted(title_fonts)[:4],
        "body_fonts":    sorted(body_fonts)[:4],
        "text_colors":   sorted(colors)[:16],
        "bg_colors":     sorted(bg_colors)[:6],
        "font_size_avg": avg_font_size,
        "font_size_max": max_font_size,
    }


# ─── KNOWLEDGE BASE LEARNER ───────────────────────────────────────────────────

def _dominant_colors(color_counts: dict, top_n: int = 8) -> list[str]:
    """Return top N colors by frequency, excluding near-white and near-black."""
    def is_neutral(hex_str):
        try:
            r, g, b = int(hex_str[0:2],16), int(hex_str[2:4],16), int(hex_str[4:6],16)
            brightness = (r + g + b) / 3
            return brightness > 230 or brightness < 20
        except Exception:
            return True

    ranked = sorted(
        [(c, n) for c, n in color_counts.items() if not is_neutral(c)],
        key=lambda x: x[1], reverse=True
    )
    return [c for c, _ in ranked[:top_n]]


def learn_design_system(reference_files: list[bytes]) -> dict:
    """
    Analyse a corpus of reference PPTXs to extract a learned design system:
    - Color palette (dominant + accent colors)
    - Font pairings (title vs body)
    - Layout patterns (which layout types appear most)
    - Spacing & sizing norms
    - Slide structure patterns (title position, content zones)
    Returns a structured dict that is injected into the AI prompt.
    """
    from collections import Counter, defaultdict

    color_counter  = Counter()
    bg_counter     = Counter()
    font_counter   = Counter()
    title_font_ctr = Counter()
    body_font_ctr  = Counter()
    font_sizes_all = []
    layout_patterns = []
    shape_counts   = []
    slide_dims     = []

    for file_bytes in reference_files:
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            w   = round(prs.slide_width  / 914400, 2)
            h   = round(prs.slide_height / 914400, 2)
            slide_dims.append((w, h))

            for slide in prs.slides:
                n_shapes   = len(slide.shapes)
                n_text     = sum(1 for s in slide.shapes if s.has_text_frame)
                n_images   = sum(1 for s in slide.shapes if s.shape_type == 13)
                n_tables   = sum(1 for s in slide.shapes if s.has_table)
                shape_counts.append(n_shapes)

                # Infer layout pattern from shape mix
                if n_images >= 1 and n_text >= 1:
                    layout_patterns.append("image+text")
                elif n_tables >= 1:
                    layout_patterns.append("table")
                elif n_text >= 3:
                    layout_patterns.append("multi-text")
                elif n_text == 2:
                    layout_patterns.append("title+body")
                elif n_text == 1:
                    layout_patterns.append("title-only")
                else:
                    layout_patterns.append("visual-only")

                # Background fill
                try:
                    bg = slide.background.fill
                    if bg.type:
                        bg_counter[str(bg.fore_color.rgb)] += 1
                except Exception:
                    pass

                for shape in slide.shapes:
                    # Shape fill (accent blocks, colored boxes)
                    try:
                        fill = shape.fill
                        if fill.type == 1:
                            color_counter[str(fill.fore_color.rgb)] += 3  # weight higher
                    except Exception:
                        pass

                    if not shape.has_text_frame:
                        continue

                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                font_counter[run.font.name] += 1
                                sz = run.font.size or 0
                                pt = round(sz / 12700)
                                font_sizes_all.append(pt)
                                if pt >= 24:
                                    title_font_ctr[run.font.name] += 1
                                elif 0 < pt <= 18:
                                    body_font_ctr[run.font.name] += 1
                            if run.font.color and run.font.color.type:
                                try:
                                    color_counter[str(run.font.color.rgb)] += 1
                                except Exception:
                                    pass

        except Exception as e:
            continue  # skip corrupt files silently

    # ── Summarise findings ────────────────────────────────────────────────────

    layout_ctr = Counter(layout_patterns)
    total_slides = sum(layout_ctr.values()) or 1

    dominant_accent_colors = _dominant_colors(color_counter, top_n=8)
    dominant_bg_colors     = _dominant_colors(bg_counter,    top_n=4)

    top_title_fonts = [f for f, _ in title_font_ctr.most_common(3)]
    top_body_fonts  = [f for f, _ in body_font_ctr.most_common(3)]
    top_all_fonts   = [f for f, _ in font_counter.most_common(6)]

    avg_font_pt  = round(sum(font_sizes_all) / len(font_sizes_all)) if font_sizes_all else 0
    title_pt_est = max((pt for pt in font_sizes_all if pt >= 24), default=32)
    body_pt_est  = round(
        sum(pt for pt in font_sizes_all if 0 < pt <= 20) /
        max(1, sum(1 for pt in font_sizes_all if 0 < pt <= 20))
    )

    avg_shapes = round(sum(shape_counts) / len(shape_counts)) if shape_counts else 0

    # Most common slide dimensions
    dim_ctr = Counter(slide_dims)
    most_common_dim = dim_ctr.most_common(1)[0][0] if dim_ctr else (13.33, 7.5)

    layout_distribution = {
        pat: f"{round(cnt / total_slides * 100)}%"
        for pat, cnt in layout_ctr.most_common()
    }

    learned = {
        "source": "learned_from_reference_repository",
        "files_analysed": len(reference_files),
        "total_slides_analysed": total_slides,
        "slide_dimensions": {
            "width_inches":  most_common_dim[0],
            "height_inches": most_common_dim[1],
            "aspect_ratio":  "16:9" if abs(most_common_dim[0] / most_common_dim[1] - 16/9) < 0.1 else "4:3",
        },
        "color_palette": {
            "accent_colors":     dominant_accent_colors,
            "background_colors": dominant_bg_colors,
            "primary_accent":    dominant_accent_colors[0] if dominant_accent_colors else None,
            "secondary_accent":  dominant_accent_colors[1] if len(dominant_accent_colors) > 1 else None,
        },
        "typography": {
            "title_fonts":      top_title_fonts,
            "body_fonts":       top_body_fonts,
            "all_detected":     top_all_fonts,
            "recommended_pair": {
                "display": top_title_fonts[0] if top_title_fonts else "Georgia",
                "body":    top_body_fonts[0]  if top_body_fonts  else "Calibri",
            },
            "title_size_pt": title_pt_est,
            "body_size_pt":  body_pt_est,
            "avg_size_pt":   avg_font_pt,
        },
        "layout_patterns": {
            "distribution":          layout_distribution,
            "most_common_layout":    layout_ctr.most_common(1)[0][0] if layout_ctr else "title+body",
            "avg_shapes_per_slide":  avg_shapes,
            "uses_images":           layout_ctr.get("image+text", 0) > 0,
            "uses_tables":           layout_ctr.get("table", 0) > 0,
            "visual_slide_ratio":    f"{round(layout_ctr.get('image+text', 0) / total_slides * 100)}%",
        },
        "design_rules_inferred": _infer_design_rules(
            dominant_accent_colors, dominant_bg_colors,
            top_title_fonts, top_body_fonts,
            layout_ctr, total_slides,
        ),
    }
    return learned


def _infer_design_rules(accent_colors, bg_colors, title_fonts, body_fonts,
                        layout_ctr, total_slides) -> list[str]:
    """Generate plain-English design rules from the extracted patterns."""
    rules = []

    if accent_colors:
        rules.append(f"Primary brand accent color is #{accent_colors[0]} — use for titles, highlights, and key shapes")
    if len(accent_colors) > 1:
        rules.append(f"Secondary accent color is #{accent_colors[1]} — use for supporting elements")

    if bg_colors:
        rules.append(f"Most common slide background is #{bg_colors[0]}")

    if title_fonts:
        rules.append(f"Use '{title_fonts[0]}' for all slide titles and headings")
    if body_fonts:
        rules.append(f"Use '{body_fonts[0]}' for body text and bullet points")

    image_pct = round(layout_ctr.get("image+text", 0) / total_slides * 100)
    if image_pct >= 40:
        rules.append(f"Reference deck is highly visual ({image_pct}% of slides use images) — prefer image+text layouts")
    elif image_pct >= 15:
        rules.append(f"Reference deck uses images on ~{image_pct}% of slides — include visuals on key slides")

    if layout_ctr.get("multi-text", 0) / total_slides > 0.4:
        rules.append("Reference deck uses multi-column text layouts frequently — apply 2-column where relevant")

    if layout_ctr.get("title-only", 0) / total_slides > 0.1:
        rules.append("Reference deck uses section-break / title-only slides — maintain that rhythm for section dividers")

    return rules

# ─── SYSTEM PROMPT ───────────────────────────────────────────────────────────

SYSTEM_PROMPT = """\
You are a world-class presentation designer and data storytelling expert.

Your task is to transform raw PowerPoint slide content into a high-end, consulting-grade presentation.

You do NOT summarize slides. You RESTRUCTURE, REDESIGN, and ELEVATE them.

CORE PRINCIPLES:
- One idea per slide
- Insight-driven titles (not descriptive)
- Max 3–5 bullets per slide
- Each bullet ≤ 8 words
- No paragraphs, no long sentences
- Prefer visuals over text
- Preserve ALL input images (never drop them)
- Split dense slides into multiple slides
- Merge redundant slides
- Ensure consistent tone, structure, and flow

DESIGN INTELLIGENCE:
- Apply the provided design system (colors, fonts, spacing)
- Select the best layout for each slide
- Maintain visual consistency across the deck
- Reposition images intelligently (left, right, full, background)

VISUAL RULES:
- Use diagrams for flows, pipelines, architectures
- Use 2-column layouts for comparisons
- Use simple bullet layouts for key points
- Use dashboard-style layouts for metrics
- Avoid text-heavy slides at all costs

STYLE:
- Executive-level communication
- Concise, confident, sharp
- No filler words
- No repetition
- No unnecessary technical noise

STORYLINE:
- Ensure logical flow across slides:
  Problem → Solution → Architecture → Value → Risks → Summary
- Reorder slides if needed
- Remove redundancy across slides

You behave like a senior partner preparing a CTO-level presentation.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
OUTPUT FORMAT — respond ONLY with valid JSON. No markdown fences. No explanation. No preamble.
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

JSON SCHEMA:
{
  "deck_title": "string — sharp, punchy deck title",
  "deck_theme": "dark | light",
  "total_slides": number,
  "slides": [
    {
      "id": number,
      "title": "Insight-driven title — the takeaway, not a label",
      "subtitle": "optional short subtitle",
      "classification": "architecture | process | executive | comparison | metrics | technical | summary",
      "layout": "title-hero | bullet | two-column | diagram | dashboard | image-right | image-left | section-break",
      "visual_type": "bullets | diagram | metric-cards | comparison | flow | image",
      "bullets": ["each bullet ≤ 8 words", "max 5 bullets"],
      "left_column": { "heading": "string", "bullets": ["..."] },
      "right_column": { "heading": "string", "bullets": ["..."] },
      "metrics": [{ "value": "80%", "label": "Time Saved" }],
      "flow_steps": ["Step A", "→", "Step B", "→", "Step C"],
      "images": ["image_id — preserve ALL images from input"],
      "image_placement": "left | right | full | background | none",
      "notes": "1–2 sentence speaker note, exec-level framing",
      "design": {
        "theme": "light | dark",
        "accent": "high | medium | low"
      }
    }
  ]
}

LAYOUT SELECTION GUIDE:
- title-hero    → opening/closing slides, section covers
- bullet        → key points, findings, recommendations
- two-column    → comparisons, pros/cons, before/after
- diagram       → flows, pipelines, architectures, processes
- dashboard     → KPIs, success metrics, performance data
- image-right   → content with supporting image right
- image-left    → image-led storytelling with content right
- section-break → major section dividers

QUALITY CHECK BEFORE OUTPUT:
- Ensure full deck has logical flow: Problem → Solution → Architecture → Value → Risks → Summary
- Ensure no redundancy across slides
- Ensure titles are strong insights (not labels)
- Ensure all input images are referenced in output
- Ensure every slide is concise and presentation-ready
"""

# ─── LAYOUT TEMPLATES ────────────────────────────────────────────────────────

LAYOUT_TEMPLATES = [
    "title-hero",       # Large title + subtitle, minimal content
    "bullet",           # Title + 3–5 bullet points
    "two-column",       # Split: left heading+bullets / right heading+bullets
    "diagram",          # Title + flow/architecture visual description
    "dashboard",        # Title + metric cards (value + label)
    "image-right",      # Bullets left, image right
    "image-left",       # Image left, content right
    "section-break",    # Divider slide between major sections
]

CLASSIFICATION_TYPES = [
    "architecture", "process", "executive",
    "comparison", "metrics", "technical", "summary"
]

# ─── AI ENGINE ───────────────────────────────────────────────────────────────

def build_prompt(parsed: dict, file_design: dict, learned_design: dict | None = None) -> str:
    """
    Build the USER message — pure data payload only.
    All transformation instructions live in the SYSTEM_PROMPT above.
    """
    slides_summary = []
    for s in parsed["slides"]:
        entry = {
            "slide_number": s["slide_number"],
            "title": s["title"],
            "bullets": s["bullets"],
            "raw_text_blocks": s["raw_text_blocks"],
            "has_images": len(s["image_ids"]) > 0,
            "image_ids": s["image_ids"],
            "has_table": len(s["tables"]) > 0,
        }
        slides_summary.append(entry)

    image_meta = {
        img_id: {
            "slide": meta["slide"],
            "ext": meta["ext"],
            "width_emu": meta["width"],
            "height_emu": meta["height"],
        }
        for img_id, meta in parsed["images"].items()
    }

    # ── Design system block: learned KB takes priority over file-extracted ──
    if learned_design:
        design_block = f"""━━━ LEARNED DESIGN SYSTEM (from {learned_design['files_analysed']} reference files · {learned_design['total_slides_analysed']} slides analysed) ━━━
THIS IS YOUR AUTHORITATIVE STYLE GUIDE. Apply it faithfully to every slide.

Color Palette:
  Primary accent:   #{learned_design['color_palette']['primary_accent'] or 'N/A'}
  Secondary accent: #{learned_design['color_palette']['secondary_accent'] or 'N/A'}
  All accent colors: {learned_design['color_palette']['accent_colors']}
  Background colors: {learned_design['color_palette']['background_colors']}

Typography:
  Title font:  {learned_design['typography']['recommended_pair']['display']}  @ {learned_design['typography']['title_size_pt']}pt
  Body font:   {learned_design['typography']['recommended_pair']['body']}  @ {learned_design['typography']['body_size_pt']}pt
  All detected fonts: {learned_design['typography']['all_detected']}

Layout Patterns (from reference corpus):
  Distribution: {json.dumps(learned_design['layout_patterns']['distribution'])}
  Most common:  {learned_design['layout_patterns']['most_common_layout']}
  Avg shapes/slide: {learned_design['layout_patterns']['avg_shapes_per_slide']}
  Visual slide ratio: {learned_design['layout_patterns']['visual_slide_ratio']}

Inferred Design Rules (apply ALL of these):
{chr(10).join(f'  • {r}' for r in learned_design['design_rules_inferred'])}

━━━ INPUT FILE DESIGN CONTEXT (secondary reference) ━━━
{json.dumps(file_design, indent=2)}"""
    else:
        design_block = f"""━━━ DESIGN SYSTEM (extracted from input file only — no reference KB loaded) ━━━
{json.dumps(file_design, indent=2)}

NOTE: No reference knowledge base was provided. Upload reference PPTXs in Step 1 to enable learned design system enforcement."""

    prompt = f"""Transform the following PowerPoint deck into a consulting-grade, executive-ready presentation.

{design_block}

━━━ INPUT SLIDES ({parsed['slide_count']} slides) ━━━
{json.dumps(slides_summary, indent=2)}

━━━ AVAILABLE LAYOUTS ━━━
{json.dumps(LAYOUT_TEMPLATES, indent=2)}

━━━ CLASSIFICATION OPTIONS ━━━
{json.dumps(CLASSIFICATION_TYPES, indent=2)}

━━━ IMAGES (preserve ALL — reference by image_id in output) ━━━
{json.dumps(image_meta, indent=2)}

━━━ TRANSFORMATION RULES FOR THIS DECK ━━━
- Split any slide with more than 5 bullets into multiple slides
- Merge any slides covering the same idea
- Every title must be an insight, not a label
- For pipeline/flow content → use flow_steps with "→" separators
- For metrics/KPI content → use metrics array with value + label
- For comparison content → use left_column + right_column with headings
- All {len(image_meta)} image(s) MUST appear in the output
- Reorder slides: Problem → Solution → Architecture → Value → Risks → Summary
{"- Use the learned design system colors and fonts in the 'design' field of each slide" if learned_design else ""}

Return ONLY the JSON object. No markdown. No preamble. No explanation."""

    return prompt


def call_qwen(prompt: str, endpoint: str, model: str, temperature: float, max_tokens: int,
              system_prompt: str = None) -> dict:
    """Call llama.cpp OpenAI-compatible endpoint."""
    effective_system = system_prompt if system_prompt else SYSTEM_PROMPT
    payload = {
        "model": model,
        "messages": [
            {
                "role": "system",
                "content": effective_system,   # ← your full designer system prompt
            },
            {"role": "user", "content": prompt}
        ],
        "temperature": temperature,
        "max_tokens": max_tokens,
        "stream": False,
    }

    response = requests.post(
        endpoint,
        json=payload,
        headers={"Content-Type": "application/json"},
        timeout=300,
    )
    response.raise_for_status()
    data = response.json()
    raw = data["choices"][0]["message"]["content"].strip()

    # Strip markdown fences if model added them anyway
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    if raw.endswith("```"):
        raw = raw[:-3].strip()

    return json.loads(raw)


# ─── PPTX OUTPUT GENERATOR ───────────────────────────────────────────────────

def hex_to_rgb(hex_str: str):
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))


def generate_pptx(transformed: dict, images: dict) -> bytes:
    """Generate a styled .pptx from the transformed JSON."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    DARK_BG   = RGBColor(13, 15, 20)
    LIGHT_BG  = RGBColor(250, 251, 255)
    ACCENT    = RGBColor(79, 142, 247)
    WHITE     = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 34, 48)
    MUTED     = RGBColor(107, 114, 128)

    W = Inches(13.33)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # fully blank

    def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
        shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
        shape.line.fill.background()
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
            if line_width:
                shape.line.width = line_width
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, left, top, width, height,
                 font_name="Calibri", font_size=18, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return txBox

    def render_slide(slide_data: dict):
        sl = prs.slides.add_slide(blank_layout)
        layout  = slide_data.get("layout", "bullet")
        theme   = slide_data.get("design", {}).get("theme", transformed.get("deck_theme", "dark"))
        bg_col  = DARK_BG if theme == "dark" else LIGHT_BG
        txt_col = WHITE   if theme == "dark" else DARK_TEXT

        # Background
        add_rect(sl, 0, 0, W, H, fill_color=bg_col)

        # Accent bar (left edge)
        add_rect(sl, 0, 0, Inches(0.05), H, fill_color=ACCENT)

        title_text = slide_data.get("title", "")
        subtitle   = slide_data.get("subtitle", "")
        bullets    = slide_data.get("bullets", [])
        notes_text = slide_data.get("notes", "")

        MARGIN = Inches(0.65)
        CONTENT_W = W - MARGIN * 2

        # ── TITLE-HERO layout ─────────────────────────────────────────────
        if layout == "title-hero":
            add_text(sl, title_text,
                     MARGIN, Inches(2.5), CONTENT_W, Inches(1.5),
                     font_name="Georgia", font_size=44, bold=True, color=txt_col,
                     align=PP_ALIGN.CENTER)
            if subtitle:
                add_text(sl, subtitle,
                         MARGIN, Inches(4.2), CONTENT_W, Inches(0.7),
                         font_size=20, color=MUTED, align=PP_ALIGN.CENTER)
            # Decorative accent line
            add_rect(sl, W//2 - Inches(1.5), Inches(4.1), Inches(3), Inches(0.025), fill_color=ACCENT)

        # ── SECTION-BREAK layout ──────────────────────────────────────────
        elif layout == "section-break":
            add_rect(sl, 0, 0, W, H, fill_color=ACCENT)
            add_text(sl, title_text,
                     MARGIN, Inches(3.0), CONTENT_W, Inches(1.2),
                     font_name="Georgia", font_size=40, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        # ── BULLET layout ─────────────────────────────────────────────────
        elif layout == "bullet":
            add_text(sl, title_text,
                     MARGIN, Inches(0.55), CONTENT_W, Inches(0.85),
                     font_name="Georgia", font_size=30, bold=True, color=txt_col)
            add_rect(sl, MARGIN, Inches(1.48), Inches(1.2), Inches(0.03), fill_color=ACCENT)

            y = Inches(1.7)
            for b in bullets[:5]:
                # Bullet dot
                add_rect(sl, MARGIN, y + Inches(0.1), Inches(0.07), Inches(0.07), fill_color=ACCENT)
                add_text(sl, b,
                         MARGIN + Inches(0.22), y, CONTENT_W - Inches(0.3), Inches(0.45),
                         font_size=17, color=txt_col)
                y += Inches(0.62)

        # ── TWO-COLUMN layout ─────────────────────────────────────────────
        elif layout == "two-column":
            add_text(sl, title_text,
                     MARGIN, Inches(0.45), CONTENT_W, Inches(0.75),
                     font_name="Georgia", font_size=28, bold=True, color=txt_col)

            col_w = (CONTENT_W - Inches(0.3)) / 2
            left_col  = slide_data.get("left_column", {})
            right_col = slide_data.get("right_column", {})

            for col_data, col_x in [(left_col, MARGIN), (right_col, MARGIN + col_w + Inches(0.3))]:
                add_rect(sl, col_x, Inches(1.4), col_w, Inches(4.8),
                         fill_color=RGBColor(22, 25, 32) if theme == "dark" else RGBColor(240, 244, 255),
                         line_color=RGBColor(42, 47, 66), line_width=Pt(0.5))
                heading = col_data.get("heading", "")
                if heading:
                    add_text(sl, heading,
                             col_x + Inches(0.2), Inches(1.6), col_w - Inches(0.4), Inches(0.5),
                             font_size=15, bold=True, color=ACCENT)
                cy = Inches(2.25)
                for b in col_data.get("bullets", [])[:5]:
                    add_text(sl, f"• {b}",
                             col_x + Inches(0.2), cy, col_w - Inches(0.4), Inches(0.45),
                             font_size=14, color=txt_col)
                    cy += Inches(0.52)

        # ── DASHBOARD (metrics) layout ────────────────────────────────────
        elif layout == "dashboard":
            add_text(sl, title_text,
                     MARGIN, Inches(0.45), CONTENT_W, Inches(0.75),
                     font_name="Georgia", font_size=28, bold=True, color=txt_col)

            metrics = slide_data.get("metrics", [])
            if not metrics and bullets:
                metrics = [{"value": b.split()[0], "label": " ".join(b.split()[1:])} for b in bullets[:4]]

            num = min(len(metrics), 4)
            if num > 0:
                card_w = (CONTENT_W - Inches(0.2) * (num - 1)) / num
                for i, m in enumerate(metrics[:num]):
                    cx = MARGIN + i * (card_w + Inches(0.2))
                    add_rect(sl, cx, Inches(2.2), card_w, Inches(2.5),
                             fill_color=RGBColor(22, 25, 32) if theme == "dark" else RGBColor(240, 244, 255),
                             line_color=ACCENT, line_width=Pt(0.75))
                    add_text(sl, m.get("value", ""),
                             cx, Inches(2.75), card_w, Inches(0.9),
                             font_name="Georgia", font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
                    add_text(sl, m.get("label", ""),
                             cx, Inches(3.7), card_w, Inches(0.5),
                             font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

        # ── DIAGRAM / FLOW layout ─────────────────────────────────────────
        elif layout in ("diagram", "image-right", "image-left"):
            add_text(sl, title_text,
                     MARGIN, Inches(0.45), CONTENT_W, Inches(0.75),
                     font_name="Georgia", font_size=28, bold=True, color=txt_col)

            flow_steps = slide_data.get("flow_steps", [])
            if flow_steps:
                # Render flow as step boxes
                steps_only = [s for s in flow_steps if s != "→"]
                arrows      = len(steps_only) - 1
                total_w     = CONTENT_W
                step_w      = Inches(1.6)
                arrow_w     = Inches(0.4)
                used_w      = len(steps_only) * step_w + arrows * arrow_w
                start_x     = MARGIN + (total_w - used_w) / 2
                y_pos       = Inches(2.8)

                x = start_x
                for s in flow_steps:
                    if s == "→":
                        add_text(sl, "→", x, y_pos + Inches(0.15),
                                 arrow_w, Inches(0.5),
                                 font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
                        x += arrow_w
                    else:
                        add_rect(sl, x, y_pos, step_w, Inches(0.75),
                                 fill_color=RGBColor(22, 25, 32) if theme == "dark" else RGBColor(230, 238, 255),
                                 line_color=ACCENT, line_width=Pt(0.75))
                        add_text(sl, s, x, y_pos + Inches(0.12),
                                 step_w, Inches(0.55),
                                 font_size=12, color=txt_col, align=PP_ALIGN.CENTER)
                        x += step_w

            # Bullets below flow if any
            if bullets:
                y = Inches(4.0)
                for b in bullets[:3]:
                    add_rect(sl, MARGIN, y + Inches(0.1), Inches(0.07), Inches(0.07), fill_color=ACCENT)
                    add_text(sl, b,
                             MARGIN + Inches(0.22), y, CONTENT_W - Inches(0.3), Inches(0.45),
                             font_size=15, color=txt_col)
                    y += Inches(0.56)

        # ── Default fallback ──────────────────────────────────────────────
        else:
            add_text(sl, title_text,
                     MARGIN, Inches(0.55), CONTENT_W, Inches(0.85),
                     font_name="Georgia", font_size=30, bold=True, color=txt_col)
            y = Inches(1.7)
            for b in bullets[:5]:
                add_rect(sl, MARGIN, y + Inches(0.1), Inches(0.07), Inches(0.07), fill_color=ACCENT)
                add_text(sl, b,
                         MARGIN + Inches(0.22), y, CONTENT_W - Inches(0.3), Inches(0.45),
                         font_size=17, color=txt_col)
                y += Inches(0.62)

        # ── Slide number footer ───────────────────────────────────────────
        slide_num = slide_data.get("id", "")
        add_text(sl, str(slide_num),
                 W - Inches(0.65), H - Inches(0.45), Inches(0.4), Inches(0.3),
                 font_size=9, color=MUTED, align=PP_ALIGN.RIGHT)

        # Deck title footer (small)
        deck_title = transformed.get("deck_title", "")
        add_text(sl, deck_title,
                 MARGIN, H - Inches(0.45), Inches(4), Inches(0.3),
                 font_size=9, color=MUTED)

        # ── Speaker notes ─────────────────────────────────────────────────
        if notes_text:
            notes_slide = sl.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text

        # ── Embed images ──────────────────────────────────────────────────
        img_ids = slide_data.get("images", [])
        for img_id in img_ids[:1]:  # Embed first image if present
            if img_id in images:
                try:
                    img_meta = images[img_id]
                    img_bytes = base64.b64decode(img_meta["data"])
                    img_stream = io.BytesIO(img_bytes)
                    # Place image on right side
                    pic_l = W - Inches(4.2)
                    pic_t = Inches(1.6)
                    pic_w = Inches(3.5)
                    sl.shapes.add_picture(img_stream, pic_l, pic_t, width=pic_w)
                except Exception:
                    pass

    for slide_data in transformed.get("slides", []):
        render_slide(slide_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─── SESSION STATE INIT ───────────────────────────────────────────────────────

for key, default in [
    ("logs", []),
    ("parsed", None),
    ("design_system", None),
    ("learned_design", None),       # ← from reference KB
    ("kb_file_names", []),          # ← names of loaded KB files
    ("transformed", None),
    ("pptx_bytes", None),
    ("stage", "idle"),              # idle | parsed | transformed | done
    ("system_prompt", SYSTEM_PROMPT),
]:
    if key not in st.session_state:
        st.session_state[key] = default

# ─── SIDEBAR ─────────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown("""
    <div style="padding:1.5rem 0 1rem;">
        <div style="font-family:'DM Serif Display',serif;font-size:1.3rem;color:#fff;margin-bottom:0.25rem;">
            ⚡ PPT Intelligence
        </div>
        <div style="font-size:0.72rem;color:#6b7280;font-family:'DM Mono',monospace;">
            Powered by Qwen2.5 · llama.cpp
        </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="card-header">LLM ENDPOINT</div>', unsafe_allow_html=True)
    llm_endpoint = st.text_input(
        "Base URL",
        value="http://localhost:8080/v1/chat/completions",
        label_visibility="collapsed",
    )
    llm_model = st.text_input(
        "Model name",
        value="qwen2.5-coder-14b",
        label_visibility="collapsed",
    )

    st.markdown('<div class="card-header" style="margin-top:1rem;">PARAMETERS</div>', unsafe_allow_html=True)
    temperature = st.slider("Temperature", 0.0, 1.0, 0.2, 0.05)
    max_tokens  = st.slider("Max Tokens", 1000, 8000, 4096, 256)

    st.markdown("---")

    # Connection test
    if st.button("🔌 Test Connection"):
        try:
            r = requests.post(
                llm_endpoint,
                json={
                    "model": llm_model,
                    "messages": [{"role": "user", "content": "Reply OK"}],
                    "max_tokens": 5,
                },
                timeout=10,
            )
            if r.status_code == 200:
                st.success("✓ Connected")
            else:
                st.error(f"HTTP {r.status_code}")
        except Exception as e:
            st.error(f"✗ {str(e)[:60]}")

    st.markdown("---")

    # ── Editable system prompt ────────────────────────────────────────────
    with st.expander("✏️ Edit System Prompt", expanded=False):
        st.markdown(
            '<div style="font-size:0.7rem;color:#6b7280;font-family:DM Mono,monospace;'
            'margin-bottom:0.5rem;">Modify the designer persona & rules sent to Qwen2.5</div>',
            unsafe_allow_html=True,
        )
        edited_prompt = st.text_area(
            "System Prompt",
            value=st.session_state.system_prompt,
            height=320,
            label_visibility="collapsed",
        )
        col_save, col_reset = st.columns(2)
        with col_save:
            if st.button("💾 Save", use_container_width=True):
                st.session_state.system_prompt = edited_prompt
                st.success("Saved")
        with col_reset:
            if st.button("↩ Reset", use_container_width=True):
                st.session_state.system_prompt = SYSTEM_PROMPT
                st.success("Reset to default")

    st.markdown("---")

    # Status indicator
    stage_map = {
        "idle":        ("off",  "Waiting for upload"),
        "parsed":      ("warn", "PPTX parsed"),
        "transformed": ("ok",   "AI transform done"),
        "done":        ("ok",   "PPTX generated"),
    }
    dot_cls, stage_label = stage_map.get(st.session_state.stage, ("off", "idle"))
    kb_dot   = "ok"  if st.session_state.learned_design else "off"
    kb_label = f"{st.session_state.learned_design['files_analysed']} ref files · {st.session_state.learned_design['total_slides_analysed']} slides" \
               if st.session_state.learned_design else "No KB loaded"
    st.markdown(f"""
    <div class="status-bar" style="flex-direction:column;gap:0.4rem;">
        <div class="status-item">
            <div class="status-dot {dot_cls}"></div>
            <span>{stage_label}</span>
        </div>
        <div class="status-item">
            <div class="status-dot {kb_dot}"></div>
            <span style="font-size:0.68rem;">KB: {kb_label}</span>
        </div>
        <div style="font-size:0.68rem;color:#4b5563;margin-top:0.2rem;">
            Slides parsed: {len(st.session_state.parsed["slides"]) if st.session_state.parsed else 0}<br>
            Images found: {len(st.session_state.parsed["images"]) if st.session_state.parsed else 0}<br>
            Output slides: {len(st.session_state.transformed["slides"]) if st.session_state.transformed else 0}
        </div>
    </div>
    """, unsafe_allow_html=True)

    if st.button("🔄 Reset All"):
        for k in ["logs", "parsed", "design_system", "learned_design",
                  "kb_file_names", "transformed", "pptx_bytes"]:
            st.session_state[k] = [] if k in ("logs", "kb_file_names") else None
        st.session_state.stage = "idle"
        st.rerun()

# ─── MAIN ────────────────────────────────────────────────────────────────────

st.markdown("""
<div class="hero">
    <div class="hero-badge">AI PRESENTATION OPERATING SYSTEM</div>
    <div class="hero-title">Presentation Intelligence Platform</div>
    <div class="hero-sub">
        Upload any PowerPoint · AI restructures &amp; redesigns · Download executive-grade deck
    </div>
</div>
""", unsafe_allow_html=True)

# ── STEP 1: KNOWLEDGE BASE ────────────────────────────────────────────────────

st.markdown('<div class="section-label">01 · REFERENCE KNOWLEDGE BASE</div>', unsafe_allow_html=True)

col_kb, col_kb_info = st.columns([2, 1])

with col_kb:
    kb_files = st.file_uploader(
        "Upload reference PPTXs (your company's best decks)",
        type=["pptx"],
        accept_multiple_files=True,
        label_visibility="collapsed",
        key="kb_uploader",
    )

with col_kb_info:
    st.markdown("""
    <div class="card">
        <div class="card-header">WHY THIS MATTERS</div>
        <ul class="bullet-list">
            <li>Teaches Qwen your brand colors & fonts</li>
            <li>Learns your layout preferences</li>
            <li>Infers spacing & visual density rules</li>
            <li>More files = stronger design signal</li>
        </ul>
        <div style="margin-top:0.75rem;font-size:0.72rem;color:#6b7280;">
            Tip: Upload 3–10 of your best existing decks
        </div>
    </div>
    """, unsafe_allow_html=True)

# Process KB files when uploaded
if kb_files:
    new_names = sorted([f.name for f in kb_files])
    if new_names != st.session_state.kb_file_names:
        with st.spinner(f"Learning design system from {len(kb_files)} reference file(s)…"):
            try:
                kb_bytes_list = [f.read() for f in kb_files]
                learned = learn_design_system(kb_bytes_list)
                st.session_state.learned_design  = learned
                st.session_state.kb_file_names   = new_names
                log(f"KB learned from {learned['files_analysed']} files · {learned['total_slides_analysed']} slides", "ok")
                log(f"Primary accent: #{learned['color_palette']['primary_accent']}", "info")
                log(f"Title font: {learned['typography']['recommended_pair']['display']} · Body: {learned['typography']['recommended_pair']['body']}", "info")
                log(f"Layout mix: {learned['layout_patterns']['distribution']}", "info")
                for rule in learned['design_rules_inferred']:
                    log(f"Rule: {rule}", "info")
            except Exception as e:
                log(f"KB learning error: {e}", "err")
                st.error(f"❌ Failed to process reference files: {e}")

# Show learned design system summary
if st.session_state.learned_design:
    ld = st.session_state.learned_design
    cp = ld["color_palette"]
    ty = ld["typography"]
    lp = ld["layout_patterns"]

    st.markdown(f"""
    <div class="card" style="border-color:rgba(79,142,247,0.3);margin-top:0.5rem;">
        <div class="card-header">✓ LEARNED DESIGN SYSTEM — {ld['files_analysed']} files · {ld['total_slides_analysed']} slides</div>
        <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:1rem;">

            <div>
                <div style="font-size:0.72rem;color:#6b7280;margin-bottom:0.4rem;font-family:DM Mono,monospace;">COLOR PALETTE</div>
                <div style="display:flex;flex-wrap:wrap;gap:0.3rem;margin-bottom:0.5rem;">
                    {"".join(f'<div title="#{c}" style="width:22px;height:22px;border-radius:4px;background:#{c};border:1px solid rgba(255,255,255,0.1);"></div>' for c in cp["accent_colors"][:8])}
                </div>
                <div style="font-size:0.72rem;color:#9ca3af;">
                    Primary: <code style="color:#4f8ef7;">#{cp["primary_accent"] or "N/A"}</code><br>
                    Secondary: <code style="color:#4f8ef7;">#{cp["secondary_accent"] or "N/A"}</code>
                </div>
            </div>

            <div>
                <div style="font-size:0.72rem;color:#6b7280;margin-bottom:0.4rem;font-family:DM Mono,monospace;">TYPOGRAPHY</div>
                <div style="font-size:0.8rem;color:#e8eaf0;">
                    <span style="color:#4f8ef7;">Title:</span> {ty["recommended_pair"]["display"]} @ {ty["title_size_pt"]}pt<br>
                    <span style="color:#4f8ef7;">Body:</span>  {ty["recommended_pair"]["body"]} @ {ty["body_size_pt"]}pt
                </div>
            </div>

            <div>
                <div style="font-size:0.72rem;color:#6b7280;margin-bottom:0.4rem;font-family:DM Mono,monospace;">LAYOUT MIX</div>
                {"".join(f'<div style="font-size:0.72rem;color:#9ca3af;line-height:1.7;">{pat}: <span style="color:#4f8ef7;">{pct}</span></div>' for pat, pct in list(lp["distribution"].items())[:5])}
            </div>

        </div>
        <div style="margin-top:0.75rem;border-top:1px solid #2a2f42;padding-top:0.75rem;">
            <div style="font-size:0.7rem;color:#6b7280;font-family:DM Mono,monospace;margin-bottom:0.35rem;">INFERRED RULES</div>
            {"".join(f'<div style="font-size:0.75rem;color:#c8cdd8;padding:0.15rem 0;">▸ {r}</div>' for r in ld["design_rules_inferred"])}
        </div>
    </div>
    """, unsafe_allow_html=True)
else:
    st.markdown("""
    <div style="padding:0.75rem 1rem;background:rgba(245,158,11,0.08);border:1px solid rgba(245,158,11,0.2);
                border-radius:8px;font-size:0.8rem;color:#fbbf24;margin-top:0.5rem;">
        ⚠️ No reference knowledge base loaded — AI will use default design judgment.
        Upload your company's best PPTXs above for branded output.
    </div>
    """, unsafe_allow_html=True)

st.markdown("---")

# ── STEP 2: UPLOAD TARGET PPTX ───────────────────────────────────────────────

st.markdown('<div class="section-label">02 · UPLOAD TARGET PRESENTATION</div>', unsafe_allow_html=True)

col_up, col_info = st.columns([2, 1])

with col_up:
    uploaded = st.file_uploader(
        "Drop your .pptx file here",
        type=["pptx"],
        label_visibility="collapsed",
    )

with col_info:
    st.markdown("""
    <div class="card">
        <div class="card-header">WHAT HAPPENS</div>
        <ul class="bullet-list">
            <li>Parse all slides, text &amp; images</li>
            <li>Apply learned design system</li>
            <li>Send to Qwen2.5 for transformation</li>
            <li>Generate executive-grade .pptx</li>
        </ul>
    </div>
    """, unsafe_allow_html=True)

if uploaded and st.session_state.stage == "idle":
    with st.spinner("Parsing PPTX…"):
        try:
            file_bytes = uploaded.read()
            st.session_state.parsed        = parse_pptx(file_bytes)
            st.session_state.design_system = extract_design_system(file_bytes)
            st.session_state.stage         = "parsed"
            log(f"Parsed {st.session_state.parsed['slide_count']} slides", "ok")
            log(f"Found {len(st.session_state.parsed['images'])} images", "ok")
            log(f"File fonts: {st.session_state.design_system['all_fonts']}", "info")
        except Exception as e:
            log(f"Parse error: {e}", "err")
            st.error(f"Failed to parse PPTX: {e}")

# ── STEP 3: PARSE SUMMARY ─────────────────────────────────────────────────────

if st.session_state.parsed:
    parsed = st.session_state.parsed
    ds     = st.session_state.design_system

    st.markdown('<div class="section-label">03 · PARSED CONTENT</div>', unsafe_allow_html=True)

    c1, c2, c3, c4 = st.columns(4)
    for col, val, label in [
        (c1, parsed["slide_count"],       "Slides"),
        (c2, len(parsed["images"]),        "Images"),
        (c3, ds.get("aspect_ratio", "?"), "Aspect"),
        (c4, len(ds.get("all_fonts", [])), "Fonts detected"),
    ]:
        col.markdown(f"""
        <div class="metric-card" style="text-align:center;padding:1rem;">
            <div class="metric-value">{val}</div>
            <div class="metric-label">{label}</div>
        </div>
        """, unsafe_allow_html=True)

    with st.expander("📋 View parsed slides"):
        for s in parsed["slides"]:
            st.markdown(f"""
            <div class="slide-card">
                <div class="slide-num">SLIDE {s['slide_number']}</div>
                <div class="slide-title">{s['title'] or '(no title)'}</div>
                <ul class="bullet-list">
                    {"".join(f"<li>{b}</li>" for b in s['bullets'][:6])}
                    {"".join(f"<li>{b}</li>" for b in s['raw_text_blocks'][:3])}
                </ul>
                {"<div style='margin-top:0.5rem;font-size:0.72rem;color:#4f8ef7;font-family:DM Mono,monospace;'>🖼 " + ", ".join(s['image_ids']) + "</div>" if s['image_ids'] else ""}
            </div>
            """, unsafe_allow_html=True)

# ── STEP 4: AI TRANSFORM ──────────────────────────────────────────────────────

if st.session_state.stage == "parsed":
    st.markdown('<div class="section-label">04 · AI TRANSFORM</div>', unsafe_allow_html=True)

    kb_loaded = st.session_state.learned_design is not None
    if not kb_loaded:
        st.markdown("""
        <div style="padding:0.6rem 1rem;background:rgba(245,158,11,0.08);border:1px solid rgba(245,158,11,0.2);
                    border-radius:8px;font-size:0.8rem;color:#fbbf24;margin-bottom:0.75rem;">
            ⚠️ No KB loaded — transformation will use default design judgment.
            Go back to Step 01 to upload reference PPTXs.
        </div>
        """, unsafe_allow_html=True)
    else:
        ld = st.session_state.learned_design
        st.markdown(f"""
        <div style="padding:0.6rem 1rem;background:rgba(34,197,94,0.08);border:1px solid rgba(34,197,94,0.2);
                    border-radius:8px;font-size:0.8rem;color:#4ade80;margin-bottom:0.75rem;">
            ✓ Design KB active — {ld['files_analysed']} ref files · {ld['total_slides_analysed']} slides ·
            Primary accent: #{ld['color_palette']['primary_accent']} ·
            Fonts: {ld['typography']['recommended_pair']['display']} / {ld['typography']['recommended_pair']['body']}
        </div>
        """, unsafe_allow_html=True)

    if st.button("⚡ Transform with Qwen2.5", use_container_width=True):
        prompt = build_prompt(
            st.session_state.parsed,
            st.session_state.design_system,
            learned_design=st.session_state.learned_design,   # ← KB injected here
        )
        active_sysprompt = st.session_state.system_prompt
        log(f"System prompt: {len(active_sysprompt):,} chars | User prompt: {len(prompt):,} chars", "info")
        log(f"Sending to {llm_endpoint}", "info")

        progress = st.progress(0, text="Calling Qwen2.5…")

        try:
            progress.progress(20, text="Waiting for model response…")
            result = call_qwen(
                prompt, llm_endpoint, llm_model, temperature, max_tokens,
                system_prompt=active_sysprompt,
            )
            progress.progress(80, text="Parsing JSON response…")
            st.session_state.transformed = result
            st.session_state.stage = "transformed"
            log(f"Transformation complete — {len(result.get('slides', []))} output slides", "ok")
            log(f"Deck title: {result.get('deck_title', '')}", "info")
            progress.progress(100, text="Done!")
            st.rerun()
        except requests.exceptions.ConnectionError:
            log("Cannot reach llama.cpp server. Is it running?", "err")
            st.error("❌ Cannot connect to llama.cpp server. Check the endpoint in the sidebar.")
            progress.empty()
        except json.JSONDecodeError as e:
            log(f"JSON parse error: {e}", "err")
            st.error("❌ Model returned invalid JSON. Try lowering temperature or increasing max_tokens.")
            progress.empty()
        except Exception as e:
            log(f"Error: {e}", "err")
            st.error(f"❌ {e}")
            progress.empty()

# ── STEP 5: RESULTS ───────────────────────────────────────────────────────────

if st.session_state.transformed:
    t = st.session_state.transformed

    st.markdown('<div class="section-label">05 · TRANSFORMED DECK</div>', unsafe_allow_html=True)

    st.markdown(f"""
    <div class="deck-header">
        <div class="deck-title">{t.get('deck_title', 'Transformed Deck')}</div>
        <div class="deck-meta">
            {len(t.get('slides', []))} slides · Theme: {t.get('deck_theme','dark')} · 
            Original: {st.session_state.parsed['slide_count']} slides
        </div>
    </div>
    """, unsafe_allow_html=True)

    # Show transformed slides
    for slide in t.get("slides", []):
        layout     = slide.get("layout", "bullet")
        cls_tag    = slide.get("classification", "")
        visual_tag = slide.get("visual_type", "")
        bullets    = slide.get("bullets", [])
        flow_steps = slide.get("flow_steps", [])
        metrics    = slide.get("metrics", [])
        left_col   = slide.get("left_column", {})
        right_col  = slide.get("right_column", {})

        # Build bullets HTML
        bullets_html = "".join(f"<li>{b}</li>" for b in bullets)

        # Build flow HTML
        flow_html = ""
        if flow_steps:
            items = []
            for step in flow_steps:
                if step == "→":
                    items.append('<span class="flow-arrow">→</span>')
                else:
                    items.append(f'<span class="flow-step">{step}</span>')
            flow_html = f'<div class="flow-steps">{"".join(items)}</div>'

        # Build metrics HTML
        metrics_html = ""
        if metrics:
            cards = "".join(
                f'<div class="metric-card"><div class="metric-value">{m.get("value","")}</div>'
                f'<div class="metric-label">{m.get("label","")}</div></div>'
                for m in metrics
            )
            metrics_html = f'<div class="metric-grid">{cards}</div>'

        # Build two-column HTML
        two_col_html = ""
        if layout == "two-column" and (left_col or right_col):
            lh = left_col.get("heading", "")
            lb = "".join(f"<li>{b}</li>" for b in left_col.get("bullets", []))
            rh = right_col.get("heading", "")
            rb = "".join(f"<li>{b}</li>" for b in right_col.get("bullets", []))
            two_col_html = f"""
            <div class="two-col">
                <div class="col-box">
                    <div class="col-heading">{lh}</div>
                    <ul class="bullet-list">{lb}</ul>
                </div>
                <div class="col-box">
                    <div class="col-heading">{rh}</div>
                    <ul class="bullet-list">{rb}</ul>
                </div>
            </div>"""

        notes_html = ""
        if slide.get("notes"):
            notes_html = f'<div style="margin-top:0.75rem;font-size:0.75rem;color:#6b7280;font-style:italic;border-left:2px solid #2a2f42;padding-left:0.75rem;">{slide["notes"]}</div>'

        img_badge = ""
        if slide.get("images"):
            img_badge = f'<span class="tag tag-visual">🖼 {len(slide["images"])} image(s)</span>'

        st.markdown(f"""
        <div class="slide-card">
            <div class="slide-num">SLIDE {slide.get('id', '')} · {layout.upper()}</div>
            <div class="slide-title">{slide.get('title', '')}</div>
            {"<div style='font-size:0.82rem;color:#9ca3af;margin-bottom:0.5rem;'>" + slide.get('subtitle','') + "</div>" if slide.get('subtitle') else ''}
            <div class="slide-meta">
                <span class="tag tag-class">{cls_tag}</span>
                <span class="tag tag-layout">{layout}</span>
                <span class="tag tag-visual">{visual_tag}</span>
                {img_badge}
            </div>
            {"<ul class='bullet-list'>" + bullets_html + "</ul>" if bullets_html else ""}
            {flow_html}
            {metrics_html}
            {two_col_html}
            {notes_html}
        </div>
        """, unsafe_allow_html=True)

    # ── Raw JSON expander
    with st.expander("🔍 View raw JSON output"):
        st.json(t)

# ── STEP 6: GENERATE PPTX ─────────────────────────────────────────────────────

if st.session_state.transformed and st.session_state.stage == "transformed":
    st.markdown('<div class="section-label">06 · GENERATE & DOWNLOAD</div>', unsafe_allow_html=True)

    col_gen, col_dl = st.columns([1, 1])

    with col_gen:
        if st.button("🏗️ Generate .pptx File", use_container_width=True):
            with st.spinner("Building PowerPoint…"):
                try:
                    pptx_bytes = generate_pptx(
                        st.session_state.transformed,
                        st.session_state.parsed["images"],
                    )
                    st.session_state.pptx_bytes = pptx_bytes
                    st.session_state.stage = "done"
                    log(f"PPTX generated — {len(pptx_bytes):,} bytes", "ok")
                    st.rerun()
                except Exception as e:
                    log(f"PPTX generation error: {e}", "err")
                    st.error(f"❌ {e}\n\n{traceback.format_exc()}")

    with col_dl:
        if st.session_state.pptx_bytes:
            deck_title = st.session_state.transformed.get("deck_title", "transformed_deck")
            safe_name  = "".join(c if c.isalnum() or c in " -_" else "_" for c in deck_title)
            st.download_button(
                label="⬇️ Download Executive Deck",
                data=st.session_state.pptx_bytes,
                file_name=f"{safe_name}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

# ── LOG PANEL ─────────────────────────────────────────────────────────────────

if st.session_state.logs:
    st.markdown('<div class="section-label">PIPELINE LOG</div>', unsafe_allow_html=True)
    render_logs()
