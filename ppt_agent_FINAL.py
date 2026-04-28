"""
╔══════════════════════════════════════════════════════════════════════╗
║   AI PRESENTATION INTELLIGENCE PLATFORM  v3.0-FINAL                ║
║   Powered by Qwen2.5 via llama.cpp  (n_ctx=8192)                   ║
╚══════════════════════════════════════════════════════════════════════╝

Install:
    pip install streamlit python-pptx requests pillow lxml

Run:
    streamlit run ppt_agent_FINAL.py

Requires llama.cpp server running at http://localhost:8080
"""

import io
import json
import base64
import requests
import traceback
from collections import Counter

import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PAGE CONFIG ─────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="PPT Intelligence Platform",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─── CSS ─────────────────────────────────────────────────────────────────────

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Serif+Display&family=DM+Mono:wght@400;500&family=DM+Sans:wght@300;400;500;600&display=swap');
:root{--bg:#0d0f14;--surface:#161920;--surface2:#1e2230;--border:#2a2f42;
      --accent:#4f8ef7;--success:#22c55e;--warning:#f59e0b;--danger:#ef4444;
      --text:#e8eaf0;--muted:#6b7280;
      --font-d:'DM Serif Display',serif;--font-b:'DM Sans',sans-serif;--font-m:'DM Mono',monospace;}
html,body,[class*="css"]{font-family:var(--font-b)!important;background:var(--bg)!important;color:var(--text)!important;}
#MainMenu,footer,header{visibility:hidden;}
.block-container{padding:2rem 2rem 4rem!important;max-width:1400px!important;}
[data-testid="stSidebar"]{background:var(--surface)!important;border-right:1px solid var(--border)!important;}
[data-testid="stSidebar"] *{color:var(--text)!important;}

.hero{background:linear-gradient(135deg,#0d0f14 0%,#161b2e 50%,#0d0f14 100%);
      border:1px solid var(--border);border-radius:16px;padding:2.5rem 3rem;
      margin-bottom:2rem;position:relative;overflow:hidden;}
.hero::before{content:'';position:absolute;top:-50%;left:-20%;width:60%;height:200%;
      background:radial-gradient(ellipse,rgba(79,142,247,.08) 0%,transparent 70%);pointer-events:none;}
.hero-badge{display:inline-block;background:rgba(79,142,247,.12);border:1px solid rgba(79,142,247,.3);
      color:var(--accent);padding:.2rem .75rem;border-radius:100px;font-size:.75rem;
      font-family:var(--font-m);margin-bottom:1rem;letter-spacing:.05em;}
.hero-title{font-family:var(--font-d)!important;font-size:2.4rem!important;font-weight:400!important;
      color:#fff!important;margin:0 0 .4rem!important;line-height:1.1!important;}
.hero-sub{font-size:.95rem;color:var(--muted);font-weight:300;letter-spacing:.02em;}

.card{background:var(--surface);border:1px solid var(--border);border-radius:12px;padding:1.5rem;margin-bottom:1rem;}
.card-header{font-size:.7rem;font-family:var(--font-m);color:var(--muted);letter-spacing:.1em;text-transform:uppercase;margin-bottom:.75rem;}

.slide-card{background:var(--surface);border:1px solid var(--border);border-radius:12px;
      padding:1.5rem 1.75rem;margin-bottom:1rem;position:relative;overflow:hidden;}
.slide-card::before{content:'';position:absolute;left:0;top:0;bottom:0;width:3px;
      background:var(--accent);border-radius:3px 0 0 3px;}
.slide-num{font-family:var(--font-m);font-size:.7rem;color:var(--accent);letter-spacing:.1em;margin-bottom:.4rem;}
.slide-title{font-family:var(--font-d)!important;font-size:1.15rem!important;color:#fff!important;
      margin-bottom:.6rem!important;line-height:1.3!important;}
.slide-meta{display:flex;gap:.5rem;flex-wrap:wrap;margin-bottom:.75rem;}
.tag{font-size:.68rem;font-family:var(--font-m);padding:.15rem .6rem;border-radius:100px;letter-spacing:.05em;}
.tag-layout{background:rgba(124,58,237,.15);border:1px solid rgba(124,58,237,.3);color:#a78bfa;}
.tag-class{background:rgba(79,142,247,.12);border:1px solid rgba(79,142,247,.3);color:var(--accent);}
.tag-visual{background:rgba(34,197,94,.1);border:1px solid rgba(34,197,94,.25);color:#4ade80;}

.bullet-list{margin:0;padding:0;list-style:none;}
.bullet-list li{font-size:.875rem;color:#c8cdd8;padding:.25rem 0 .25rem 1.1rem;position:relative;line-height:1.5;}
.bullet-list li::before{content:'▸';position:absolute;left:0;color:var(--accent);font-size:.65rem;top:.35rem;}

.flow-steps{display:flex;flex-wrap:wrap;align-items:center;gap:.4rem;margin-top:.75rem;}
.flow-step{background:var(--surface2);border:1px solid var(--border);border-radius:6px;
      padding:.3rem .75rem;font-size:.78rem;font-family:var(--font-m);color:var(--text);}
.flow-arrow{color:var(--accent);font-size:.9rem;}

.metric-grid{display:flex;gap:1rem;flex-wrap:wrap;margin-top:.75rem;}
.metric-card{background:var(--surface2);border:1px solid var(--border);border-radius:8px;
      padding:.75rem 1.25rem;text-align:center;min-width:90px;}
.metric-value{font-family:var(--font-d);font-size:1.6rem;color:var(--accent);line-height:1;}
.metric-label{font-size:.68rem;color:var(--muted);margin-top:.2rem;}

.two-col{display:grid;grid-template-columns:1fr 1fr;gap:1rem;margin-top:.75rem;}
.col-box{background:var(--surface2);border:1px solid var(--border);border-radius:8px;padding:.75rem 1rem;}
.col-heading{font-size:.78rem;font-weight:600;color:#fff;margin-bottom:.4rem;}

.status-bar{background:var(--surface);border:1px solid var(--border);border-radius:8px;
      padding:.6rem 1rem;font-family:var(--font-m);font-size:.75rem;color:var(--muted);
      margin-bottom:1rem;display:flex;gap:1.5rem;}
.status-item{display:flex;align-items:center;gap:.4rem;}
.status-dot{width:6px;height:6px;border-radius:50%;background:var(--success);}
.status-dot.warn{background:var(--warning);}
.status-dot.off{background:var(--muted);}

.log-box{background:#080a0f;border:1px solid var(--border);border-radius:8px;padding:1rem;
      font-family:var(--font-m);font-size:.72rem;color:#6b7280;max-height:200px;overflow-y:auto;line-height:1.7;}
.log-ok{color:#22c55e;}.log-info{color:#4f8ef7;}.log-warn{color:#f59e0b;}.log-err{color:#ef4444;}

.stButton>button{background:var(--accent)!important;color:#fff!important;border:none!important;
      border-radius:8px!important;font-family:var(--font-b)!important;font-weight:500!important;
      padding:.55rem 1.5rem!important;transition:opacity .2s!important;}
.stButton>button:hover{opacity:.85!important;}
.stTextInput>div>div>input,.stNumberInput>div>div>input{background:var(--surface2)!important;
      border:1px solid var(--border)!important;border-radius:8px!important;
      color:var(--text)!important;font-family:var(--font-m)!important;font-size:.85rem!important;}
.streamlit-expanderHeader{background:var(--surface)!important;border:1px solid var(--border)!important;
      border-radius:8px!important;color:var(--text)!important;}
.stProgress>div>div>div{background:var(--accent)!important;border-radius:4px!important;}
hr{border-color:var(--border)!important;margin:1.5rem 0!important;}

.deck-header{background:linear-gradient(90deg,rgba(79,142,247,.1),transparent);
      border:1px solid rgba(79,142,247,.2);border-left:3px solid var(--accent);
      border-radius:8px;padding:1rem 1.5rem;margin-bottom:1.5rem;}
.deck-title{font-family:var(--font-d);font-size:1.4rem;color:#fff;}
.deck-meta{font-size:.78rem;color:var(--muted);margin-top:.25rem;font-family:var(--font-m);}

.section-label{font-size:.65rem;font-family:var(--font-m);letter-spacing:.15em;color:var(--muted);
      text-transform:uppercase;margin:1.5rem 0 .75rem;display:flex;align-items:center;gap:.5rem;}
.section-label::after{content:'';flex:1;height:1px;background:var(--border);}
</style>
""", unsafe_allow_html=True)


# ─── LOGGING ─────────────────────────────────────────────────────────────────

def log(msg, kind="info"):
    cls    = {"ok":"log-ok","info":"log-info","warn":"log-warn","err":"log-err"}.get(kind,"log-info")
    prefix = {"ok":"✓","info":"›","warn":"⚠","err":"✗"}.get(kind,"›")
    st.session_state.logs.append('<span class="{}">{} {}</span>'.format(cls, prefix, msg))

def render_logs():
    if st.session_state.logs:
        html = "<br>".join(st.session_state.logs[-40:])
        st.markdown('<div class="log-box">{}</div>'.format(html), unsafe_allow_html=True)


# ─── SAFE SHAPE HELPERS ───────────────────────────────────────────────────────
# Every helper catches ALL exceptions — no shape attribute access ever crashes.

def _safe_placeholder_idx(shape):
    """Return placeholder idx (int) or None. Handles 'Shape is not a placeholder'."""
    try:
        pf = shape.placeholder_format
        if pf is None:
            return None
        return pf.idx          # this is the line that used to raise
    except Exception:
        return None

def _safe_text(shape):
    try:
        return shape.text_frame.text.strip()
    except Exception:
        return ""

def _safe_paragraphs(shape):
    try:
        return [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
    except Exception:
        return []

def _safe_has_text(shape):
    try:
        return bool(shape.has_text_frame)
    except Exception:
        return False

def _safe_has_table(shape):
    try:
        return bool(shape.has_table)
    except Exception:
        return False

def _safe_shape_type(shape):
    try:
        return shape.shape_type
    except Exception:
        return None


# ─── PPTX PARSER ─────────────────────────────────────────────────────────────

def parse_pptx(file_bytes):
    prs        = Presentation(io.BytesIO(file_bytes))
    slides_out = []
    images     = {}

    for slide_idx, slide in enumerate(prs.slides):
        info = {
            "slide_number":    slide_idx + 1,
            "title":           "",
            "bullets":         [],
            "tables":          [],
            "image_ids":       [],
            "raw_text_blocks": [],
        }

        try:
            shapes = list(slide.shapes)
        except Exception:
            slides_out.append(info)
            continue

        for shape in shapes:
            stype = _safe_shape_type(shape)

            # ── text ──────────────────────────────────────────────────────
            if _safe_has_text(shape):
                text = _safe_text(shape)
                if text:
                    ph = _safe_placeholder_idx(shape)   # SAFE — never raises
                    if ph == 0:
                        info["title"] = text
                    elif ph == 1:
                        info["bullets"].extend(_safe_paragraphs(shape))
                    elif ph is not None:
                        info["raw_text_blocks"].append(text)
                    else:
                        lines = _safe_paragraphs(shape)
                        if lines:
                            if not info["title"] and len(lines[0]) < 120:
                                info["title"] = lines[0]
                                info["raw_text_blocks"].extend(lines[1:])
                            else:
                                info["raw_text_blocks"].extend(lines)

            # ── table ─────────────────────────────────────────────────────
            if _safe_has_table(shape):
                try:
                    rows = [[c.text.strip() for c in r.cells] for r in shape.table.rows]
                    if rows:
                        info["tables"].append(rows)
                except Exception:
                    pass

            # ── picture ───────────────────────────────────────────────────
            if stype == 13:
                try:
                    iid = "img_{}_{}".format(slide_idx, shape.shape_id)
                    images[iid] = {
                        "data":   base64.b64encode(shape.image.blob).decode(),
                        "ext":    shape.image.ext,
                        "slide":  slide_idx + 1,
                        "width":  shape.width,
                        "height": shape.height,
                    }
                    info["image_ids"].append(iid)
                except Exception:
                    pass

            # ── group — recurse one level ──────────────────────────────────
            if stype == 6:
                try:
                    for child in shape.shapes:
                        if _safe_has_text(child):
                            info["raw_text_blocks"].extend(_safe_paragraphs(child))
                        if _safe_shape_type(child) == 13:
                            try:
                                iid = "img_{}_{}".format(slide_idx, child.shape_id)
                                images[iid] = {
                                    "data":   base64.b64encode(child.image.blob).decode(),
                                    "ext":    child.image.ext,
                                    "slide":  slide_idx + 1,
                                    "width":  child.width,
                                    "height": child.height,
                                }
                                info["image_ids"].append(iid)
                            except Exception:
                                pass
                except Exception:
                    pass

        # fallback title
        if not info["title"] and info["raw_text_blocks"]:
            info["title"] = info["raw_text_blocks"].pop(0)

        slides_out.append(info)

    return {"slide_count": len(slides_out), "slides": slides_out, "images": images}


# ─── DESIGN EXTRACTOR (single file) ──────────────────────────────────────────

def extract_design_system(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    colors, fonts, title_fonts, body_fonts, bg_colors = set(), set(), set(), set(), set()
    sizes = []

    for slide in prs.slides:
        try:
            bg = slide.background.fill
            if bg.type:
                bg_colors.add(str(bg.fore_color.rgb))
        except Exception:
            pass

        try:
            shapes = list(slide.shapes)
        except Exception:
            continue

        for shape in shapes:
            try:
                if shape.fill.type == 1:
                    colors.add(str(shape.fill.fore_color.rgb))
            except Exception:
                pass

            if not _safe_has_text(shape):
                continue
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.name:
                                fonts.add(run.font.name)
                                sz = run.font.size or 0
                                pt = round(sz / 12700)
                                sizes.append(pt)
                                if pt >= 24:
                                    title_fonts.add(run.font.name)
                                elif 0 < pt <= 18:
                                    body_fonts.add(run.font.name)
                        except Exception:
                            pass
                        try:
                            if run.font.color and run.font.color.type:
                                colors.add(str(run.font.color.rgb))
                        except Exception:
                            pass
            except Exception:
                pass

    w, h   = prs.slide_width, prs.slide_height
    aspect = "16:9" if abs(w / h - 16/9) < 0.1 else "4:3"
    return {
        "aspect_ratio":       aspect,
        "slide_width_inches": round(w / 914400, 2),
        "slide_height_in":    round(h / 914400, 2),
        "slide_count":        len(prs.slides),
        "all_fonts":          sorted(fonts)[:10],
        "title_fonts":        sorted(title_fonts)[:4],
        "body_fonts":         sorted(body_fonts)[:4],
        "text_colors":        sorted(colors)[:14],
        "bg_colors":          sorted(bg_colors)[:4],
        "font_size_avg":      round(sum(sizes)/len(sizes)) if sizes else 0,
        "font_size_max":      max(sizes) if sizes else 0,
    }


# ─── KNOWLEDGE BASE LEARNER ───────────────────────────────────────────────────

def _dominant_colors(counter, top_n=8):
    def neutral(h):
        try:
            r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
            return (r+g+b)/3 > 230 or (r+g+b)/3 < 20
        except Exception:
            return True
    ranked = sorted([(c,n) for c,n in counter.items() if not neutral(c)], key=lambda x:x[1], reverse=True)
    return [c for c,_ in ranked[:top_n]]


def learn_design_system(reference_files):
    color_ctr, bg_ctr, font_ctr = Counter(), Counter(), Counter()
    title_ctr, body_ctr         = Counter(), Counter()
    sizes_all, shape_counts, layout_patterns, slide_dims = [], [], [], []

    for file_bytes in reference_files:
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            slide_dims.append((round(prs.slide_width/914400,2), round(prs.slide_height/914400,2)))

            for slide in prs.slides:
                try:
                    all_shapes = list(slide.shapes)
                except Exception:
                    all_shapes = []

                n_text   = sum(1 for s in all_shapes if _safe_has_text(s))
                n_images = sum(1 for s in all_shapes if _safe_shape_type(s) == 13)
                n_tables = sum(1 for s in all_shapes if _safe_has_table(s))
                shape_counts.append(len(all_shapes))

                if n_images >= 1 and n_text >= 1:
                    layout_patterns.append("image+text")
                elif n_tables >= 1:
                    layout_patterns.append("table")
                elif n_text >= 3:
                    layout_patterns.append("multi-text")
                elif n_text == 2:
                    layout_patterns.append("title+body")
                elif n_text == 1:
                    layout_patterns.append("title-only")
                else:
                    layout_patterns.append("visual-only")

                try:
                    bg = slide.background.fill
                    if bg.type:
                        bg_ctr[str(bg.fore_color.rgb)] += 1
                except Exception:
                    pass

                for shape in all_shapes:
                    try:
                        if shape.fill.type == 1:
                            color_ctr[str(shape.fill.fore_color.rgb)] += 3
                    except Exception:
                        pass

                    if not _safe_has_text(shape):
                        continue
                    try:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.name:
                                        font_ctr[run.font.name] += 1
                                        pt = round((run.font.size or 0)/12700)
                                        sizes_all.append(pt)
                                        if pt >= 24:
                                            title_ctr[run.font.name] += 1
                                        elif 0 < pt <= 18:
                                            body_ctr[run.font.name] += 1
                                except Exception:
                                    pass
                                try:
                                    if run.font.color and run.font.color.type:
                                        color_ctr[str(run.font.color.rgb)] += 1
                                except Exception:
                                    pass
                    except Exception:
                        pass
        except Exception:
            continue

    layout_ctr   = Counter(layout_patterns)
    total_slides = sum(layout_ctr.values()) or 1
    accent_cols  = _dominant_colors(color_ctr, 8)
    bg_cols      = _dominant_colors(bg_ctr, 4)
    top_title    = [f for f,_ in title_ctr.most_common(3)]
    top_body     = [f for f,_ in body_ctr.most_common(3)]
    top_all      = [f for f,_ in font_ctr.most_common(6)]
    title_pt     = max((p for p in sizes_all if p >= 24), default=32)
    body_pt      = round(sum(p for p in sizes_all if 0<p<=20) / max(1, sum(1 for p in sizes_all if 0<p<=20)))
    avg_shapes   = round(sum(shape_counts)/len(shape_counts)) if shape_counts else 0
    dim_ctr      = Counter(slide_dims)
    best_dim     = dim_ctr.most_common(1)[0][0] if dim_ctr else (13.33, 7.5)
    layout_dist  = {p: "{}%".format(round(c/total_slides*100)) for p,c in layout_ctr.most_common()}

    rules = []
    if accent_cols:
        rules.append("Primary brand accent is #{} — use for titles, highlights, key shapes".format(accent_cols[0]))
    if len(accent_cols) > 1:
        rules.append("Secondary accent is #{} — use for supporting elements".format(accent_cols[1]))
    if bg_cols:
        rules.append("Most common background is #{}".format(bg_cols[0]))
    if top_title:
        rules.append("Use '{}' for all slide titles".format(top_title[0]))
    if top_body:
        rules.append("Use '{}' for all body text".format(top_body[0]))
    img_pct = round(layout_ctr.get("image+text",0)/total_slides*100)
    if img_pct >= 30:
        rules.append("Deck is highly visual ({}% image slides) — prefer image+text layouts".format(img_pct))

    return {
        "source":               "learned_from_reference_repository",
        "files_analysed":       len(reference_files),
        "total_slides_analysed": total_slides,
        "slide_dimensions":     {"width_inches": best_dim[0], "height_inches": best_dim[1]},
        "color_palette": {
            "accent_colors":     accent_cols,
            "background_colors": bg_cols,
            "primary_accent":    accent_cols[0] if accent_cols else None,
            "secondary_accent":  accent_cols[1] if len(accent_cols) > 1 else None,
        },
        "typography": {
            "title_fonts":       top_title,
            "body_fonts":        top_body,
            "all_detected":      top_all,
            "recommended_pair":  {
                "display": top_title[0] if top_title else "Georgia",
                "body":    top_body[0]  if top_body  else "Calibri",
            },
            "title_size_pt": title_pt,
            "body_size_pt":  body_pt,
        },
        "layout_patterns": {
            "distribution":         layout_dist,
            "most_common_layout":   layout_ctr.most_common(1)[0][0] if layout_ctr else "title+body",
            "avg_shapes_per_slide": avg_shapes,
            "visual_slide_ratio":   "{}%".format(img_pct),
        },
        "design_rules_inferred": rules,
    }


# ─── SYSTEM PROMPT ────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """\
You are a world-class presentation designer and data storytelling expert.

Your task is to transform raw PowerPoint slide content into a high-end, consulting-grade presentation.

You do NOT summarize slides. You RESTRUCTURE, REDESIGN, and ELEVATE them.

CORE PRINCIPLES:
- One idea per slide
- Insight-driven titles (not descriptive)
- Max 3-5 bullets per slide
- Each bullet <= 8 words
- No paragraphs, no long sentences
- Prefer visuals over text
- Preserve ALL input images (never drop them)
- Split dense slides into multiple slides
- Merge redundant slides
- Ensure consistent tone, structure, and flow

DESIGN INTELLIGENCE:
- Apply the provided design system (colors, fonts, spacing)
- Select the best layout for each slide
- Maintain visual consistency across the deck
- Reposition images intelligently (left, right, full, background)

VISUAL RULES:
- Use diagrams for flows, pipelines, architectures
- Use 2-column layouts for comparisons
- Use simple bullet layouts for key points
- Use dashboard-style layouts for metrics
- Avoid text-heavy slides at all costs

STYLE:
- Executive-level communication
- Concise, confident, sharp
- No filler words, no repetition

STORYLINE:
- Logical flow: Problem -> Solution -> Architecture -> Value -> Risks -> Summary
- Reorder slides if needed

You behave like a senior partner preparing a CTO-level presentation.

OUTPUT: Respond ONLY with valid JSON. No markdown fences. No explanation. No preamble.

JSON schema:
{
  "deck_title": "string",
  "deck_theme": "dark|light",
  "total_slides": N,
  "slides": [
    {
      "id": N,
      "title": "Insight-driven title",
      "subtitle": "optional",
      "classification": "architecture|process|executive|comparison|metrics|technical|summary",
      "layout": "title-hero|bullet|two-column|diagram|dashboard|image-right|image-left|section-break",
      "visual_type": "bullets|diagram|metric-cards|comparison|flow|image",
      "bullets": ["<=8 words each, max 5"],
      "left_column":  {"heading": "", "bullets": []},
      "right_column": {"heading": "", "bullets": []},
      "metrics":    [{"value": "", "label": ""}],
      "flow_steps": ["Step A", "->", "Step B"],
      "images":  ["image_id"],
      "notes":   "1-2 sentence speaker note",
      "design":  {"theme": "dark|light", "accent": "high|medium|low"}
    }
  ]
}
"""

LAYOUT_TEMPLATES    = ["title-hero","bullet","two-column","diagram","dashboard","image-right","image-left","section-break"]
CLASSIFICATION_TYPES = ["architecture","process","executive","comparison","metrics","technical","summary"]


# ─── PROMPT BUILDER ───────────────────────────────────────────────────────────

def _slim_slides(slides):
    out = []
    for s in slides:
        out.append({
            "n":       s["slide_number"],
            "title":   s["title"][:120],
            "bullets": [b[:100] for b in s["bullets"][:5]],
            "text":    [t[:100] for t in s["raw_text_blocks"][:3]],
            "has_img": bool(s["image_ids"]),
            "img_ids": s["image_ids"],
            "has_tbl": bool(s["tables"]),
        })
    return out


def _design_block(file_design, learned_design):
    if learned_design:
        cp = learned_design["color_palette"]
        ty = learned_design["typography"]
        lp = learned_design["layout_patterns"]
        rules = "; ".join(learned_design["design_rules_inferred"])
        return (
            "DESIGN SYSTEM (learned from {} ref files, {} slides):\n"
            "Colors: primary=#{} secondary=#{} palette={}\n"
            "Fonts: title={}@{}pt body={}@{}pt\n"
            "Layouts: {} dominant, distribution={}\n"
            "Rules: {}"
        ).format(
            learned_design["files_analysed"],
            learned_design["total_slides_analysed"],
            cp["primary_accent"] or "N/A",
            cp["secondary_accent"] or "N/A",
            cp["accent_colors"][:4],
            ty["recommended_pair"]["display"], ty["title_size_pt"],
            ty["recommended_pair"]["body"],    ty["body_size_pt"],
            lp["most_common_layout"],
            list(lp["distribution"].items())[:4],
            rules,
        )
    fonts  = file_design.get("all_fonts", [])[:4]
    colors = file_design.get("text_colors", [])[:5]
    return "DESIGN SYSTEM (file only): fonts={} colors={} aspect={}".format(
        fonts, colors, file_design.get("aspect_ratio","16:9")
    )


def build_prompt(parsed, file_design, learned_design=None, slide_batch=None):
    slides_to_use = slide_batch if slide_batch is not None else parsed["slides"]
    slimmed       = _slim_slides(slides_to_use)
    design        = _design_block(file_design, learned_design)

    all_img_ids = []
    for s in slides_to_use:
        all_img_ids.extend(s["image_ids"])
    img_list = ", ".join(all_img_ids) if all_img_ids else "none"

    total = parsed["slide_count"]
    scope = "(slides {}-{} of {})".format(
        slimmed[0]["n"], slimmed[-1]["n"], total
    ) if slide_batch else "({} slides)".format(total)

    return (
        "Transform these slides into consulting-grade output {}.\n\n"
        "{}\n\n"
        "LAYOUTS: {}\n"
        "CLASSIFICATIONS: {}\n"
        "IMAGES (preserve all): {}\n\n"
        "SLIDES:\n{}\n\n"
        "RULES:\n"
        "- Title = insight/takeaway, not a label\n"
        "- Max 5 bullets, each <=8 words\n"
        "- Split slides with >5 bullets; merge redundant ones\n"
        "- flow_steps[] for pipelines; metrics[] for KPIs; columns for comparisons\n"
        "- All image ids must appear in output\n"
        "- Story arc: Problem->Solution->Architecture->Value->Risks->Summary\n\n"
        "Return ONLY the JSON object. No markdown. No preamble."
    ).format(
        scope, design,
        LAYOUT_TEMPLATES, CLASSIFICATION_TYPES, img_list,
        json.dumps(slimmed, separators=(",",":")),
    )


# ─── JSON REPAIR ─────────────────────────────────────────────────────────────

def _repair_json(raw):
    depth_brace   = raw.count("{") - raw.count("}")
    depth_bracket = raw.count("[") - raw.count("]")
    trimmed = raw.rstrip()
    if trimmed.endswith(","):
        trimmed = trimmed[:-1]

    # detect unclosed string
    in_str, esc = False, False
    for ch in trimmed:
        if esc:       esc = False; continue
        if ch == "\\": esc = True;  continue
        if ch == '"':  in_str = not in_str
    if in_str:
        trimmed += '"'

    trimmed += "]" * max(0, depth_bracket)
    trimmed += "}" * max(0, depth_brace)
    return trimmed


# ─── LLM CALL ────────────────────────────────────────────────────────────────

def call_qwen(prompt, endpoint, model, temperature, max_tokens, system_prompt=None):
    effective_system = system_prompt or SYSTEM_PROMPT
    payload = {
        "model":       model,
        "messages":    [
            {"role": "system", "content": effective_system},
            {"role": "user",   "content": prompt},
        ],
        "temperature": temperature,
        "max_tokens":  max_tokens,
        "stream":      False,
    }
    response = requests.post(endpoint, json=payload,
                             headers={"Content-Type": "application/json"}, timeout=300)
    response.raise_for_status()
    data          = response.json()
    finish_reason = data["choices"][0].get("finish_reason", "")
    raw           = data["choices"][0]["message"]["content"].strip()

    # strip fences
    if raw.startswith("```"):
        raw = raw.split("```", 2)[1]
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    if raw.endswith("```"):
        raw = raw[:-3].strip()

    # parse — then repair if needed
    try:
        return json.loads(raw), finish_reason
    except json.JSONDecodeError:
        pass

    repaired = _repair_json(raw)
    try:
        result = json.loads(repaired)
        result["_truncated"] = True
        return result, finish_reason
    except json.JSONDecodeError as e:
        raise json.JSONDecodeError(
            "JSON failed even after repair (finish_reason={!r}). "
            "Reduce 'Slides per batch' to 2 or 3. Original: {}".format(finish_reason, e.msg),
            e.doc, e.pos,
        )


# ─── BATCH TRANSFORM ─────────────────────────────────────────────────────────

def transform_in_batches(parsed, file_design, learned_design,
                         endpoint, model, temperature, max_tokens,
                         system_prompt, batch_size, progress_cb=None):
    slides    = parsed["slides"]
    batches   = [slides[i:i+batch_size] for i in range(0, len(slides), batch_size)]
    all_slides, deck_meta, truncated = [], {}, False

    for idx, batch in enumerate(batches):
        if progress_cb:
            pct = int(10 + 75 * idx / len(batches))
            progress_cb(pct,
                "Batch {}/{} — slides {}-{}…".format(
                    idx+1, len(batches),
                    batch[0]["slide_number"], batch[-1]["slide_number"]))

        prompt = build_prompt(parsed, file_design, learned_design, slide_batch=batch)
        log("Batch {}: {} slides, prompt ~{} chars".format(idx+1, len(batch), len(prompt)), "info")

        result, finish_reason = call_qwen(
            prompt, endpoint, model, temperature, max_tokens, system_prompt)
        log("Batch {} done — finish_reason={}".format(idx+1, finish_reason),
            "warn" if finish_reason == "length" else "ok")

        if result.get("_truncated"):
            truncated = True
            log("Batch {} was truncated — JSON auto-repaired".format(idx+1), "warn")

        if not deck_meta:
            deck_meta = {
                "deck_title": result.get("deck_title", "Transformed Deck"),
                "deck_theme": result.get("deck_theme", "dark"),
            }

        for i, sl in enumerate(result.get("slides", [])):
            sl["id"] = len(all_slides) + i + 1
        all_slides.extend(result.get("slides", []))

    return {**deck_meta, "total_slides": len(all_slides),
            "slides": all_slides, "_truncated": truncated}


# ─── PPTX GENERATOR ──────────────────────────────────────────────────────────

def generate_pptx(transformed, images):
    DARK_BG   = RGBColor(13,  15,  20)
    LIGHT_BG  = RGBColor(250, 251, 255)
    ACCENT    = RGBColor(79,  142, 247)
    WHITE     = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30,  34,  48)
    MUTED     = RGBColor(107, 114, 128)
    W, H      = Inches(13.33), Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    def rect(sl, l,t,w,h, fill=None, line=None, lw=None):
        sh = sl.shapes.add_shape(1,l,t,w,h)
        sh.fill.background() if not fill else (sh.fill.solid(), setattr(sh.fill.fore_color,'rgb',fill))
        sh.line.fill.background()
        if line:
            sh.line.color.rgb = line
            if lw: sh.line.width = lw
        return sh

    def txt(sl, text, l,t,w,h, fname="Calibri", fsize=16, bold=False,
            color=WHITE, align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(l,t,w,h)
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name  = fname
        run.font.size  = Pt(fsize)
        run.font.bold  = bold
        run.font.color.rgb = color
        return tb

    def render(sd):
        sl     = prs.slides.add_slide(blank)
        layout = sd.get("layout","bullet")
        theme  = sd.get("design",{}).get("theme", transformed.get("deck_theme","dark"))
        bg     = DARK_BG  if theme=="dark" else LIGHT_BG
        tc     = WHITE    if theme=="dark" else DARK_TEXT
        M      = Inches(0.65)
        CW     = W - M*2

        rect(sl, 0,0,W,H, fill=bg)                      # background
        rect(sl, 0,0,Inches(0.05),H, fill=ACCENT)        # left stripe

        title   = sd.get("title","")
        bullets = sd.get("bullets",[])

        if layout == "title-hero":
            txt(sl, title, M, Inches(2.5), CW, Inches(1.5),
                fname="Georgia", fsize=44, bold=True, color=tc, align=PP_ALIGN.CENTER)
            sub = sd.get("subtitle","")
            if sub:
                txt(sl, sub, M, Inches(4.2), CW, Inches(0.7),
                    fsize=20, color=MUTED, align=PP_ALIGN.CENTER)

        elif layout == "section-break":
            rect(sl, 0,0,W,H, fill=ACCENT)
            txt(sl, title, M, Inches(3.0), CW, Inches(1.2),
                fname="Georgia", fsize=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        elif layout == "two-column":
            txt(sl, title, M, Inches(0.45), CW, Inches(0.75),
                fname="Georgia", fsize=28, bold=True, color=tc)
            cw = (CW - Inches(0.3))/2
            for col_data, cx in [(sd.get("left_column",{}), M),
                                  (sd.get("right_column",{}), M+cw+Inches(0.3))]:
                rect(sl, cx, Inches(1.4), cw, Inches(4.8),
                     fill=RGBColor(22,25,32) if theme=="dark" else RGBColor(240,244,255),
                     line=RGBColor(42,47,66), lw=Pt(0.5))
                if col_data.get("heading"):
                    txt(sl, col_data["heading"], cx+Inches(0.2), Inches(1.6),
                        cw-Inches(0.4), Inches(0.5), fsize=15, bold=True, color=ACCENT)
                cy = Inches(2.25)
                for b in col_data.get("bullets",[])[:5]:
                    txt(sl, "• "+b, cx+Inches(0.2), cy, cw-Inches(0.4), Inches(0.45),
                        fsize=14, color=tc)
                    cy += Inches(0.52)

        elif layout == "dashboard":
            txt(sl, title, M, Inches(0.45), CW, Inches(0.75),
                fname="Georgia", fsize=28, bold=True, color=tc)
            metrics = sd.get("metrics",[])
            if not metrics and bullets:
                metrics = [{"value": b.split()[0], "label": " ".join(b.split()[1:])} for b in bullets[:4]]
            num = min(len(metrics), 4)
            if num:
                cw = (CW - Inches(0.2)*(num-1))/num
                for i,m in enumerate(metrics[:num]):
                    cx = M + i*(cw+Inches(0.2))
                    rect(sl, cx, Inches(2.2), cw, Inches(2.5),
                         fill=RGBColor(22,25,32) if theme=="dark" else RGBColor(240,244,255),
                         line=ACCENT, lw=Pt(0.75))
                    txt(sl, m.get("value",""), cx, Inches(2.75), cw, Inches(0.9),
                        fname="Georgia", fsize=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
                    txt(sl, m.get("label",""), cx, Inches(3.7), cw, Inches(0.5),
                        fsize=13, color=MUTED, align=PP_ALIGN.CENTER)

        elif layout in ("diagram","image-right","image-left"):
            txt(sl, title, M, Inches(0.45), CW, Inches(0.75),
                fname="Georgia", fsize=28, bold=True, color=tc)
            flow = sd.get("flow_steps",[])
            if flow:
                steps_only = [s for s in flow if s not in ("->","→")]
                n = len(steps_only)
                sw, aw = Inches(1.5), Inches(0.4)
                used   = n*sw + (n-1)*aw
                x      = M + (CW-used)/2
                y      = Inches(2.6)
                for step in flow:
                    if step in ("->","→"):
                        txt(sl, "→", x, y+Inches(0.15), aw, Inches(0.5),
                            fsize=20, color=ACCENT, align=PP_ALIGN.CENTER)
                        x += aw
                    else:
                        rect(sl, x, y, sw, Inches(0.75),
                             fill=RGBColor(22,25,32) if theme=="dark" else RGBColor(230,238,255),
                             line=ACCENT, lw=Pt(0.75))
                        txt(sl, step, x, y+Inches(0.12), sw, Inches(0.55),
                            fsize=12, color=tc, align=PP_ALIGN.CENTER)
                        x += sw
            y = Inches(4.0)
            for b in bullets[:3]:
                rect(sl, M, y+Inches(0.1), Inches(0.07), Inches(0.07), fill=ACCENT)
                txt(sl, b, M+Inches(0.22), y, CW-Inches(0.3), Inches(0.45), fsize=15, color=tc)
                y += Inches(0.56)

        else:  # bullet (default)
            txt(sl, title, M, Inches(0.55), CW, Inches(0.85),
                fname="Georgia", fsize=30, bold=True, color=tc)
            rect(sl, M, Inches(1.48), Inches(1.2), Inches(0.03), fill=ACCENT)
            y = Inches(1.7)
            for b in bullets[:5]:
                rect(sl, M, y+Inches(0.1), Inches(0.07), Inches(0.07), fill=ACCENT)
                txt(sl, b, M+Inches(0.22), y, CW-Inches(0.3), Inches(0.45), fsize=17, color=tc)
                y += Inches(0.62)

        # slide number + deck footer
        txt(sl, str(sd.get("id","")),
            W-Inches(0.65), H-Inches(0.45), Inches(0.4), Inches(0.3),
            fsize=9, color=MUTED, align=PP_ALIGN.RIGHT)
        txt(sl, transformed.get("deck_title",""),
            M, H-Inches(0.45), Inches(5), Inches(0.3), fsize=9, color=MUTED)

        # speaker notes
        if sd.get("notes"):
            try:
                sl.notes_slide.notes_text_frame.text = sd["notes"]
            except Exception:
                pass

        # embed first image
        for iid in sd.get("images",[])[:1]:
            if iid in images:
                try:
                    img_stream = io.BytesIO(base64.b64decode(images[iid]["data"]))
                    sl.shapes.add_picture(img_stream, W-Inches(4.2), Inches(1.6), width=Inches(3.5))
                except Exception:
                    pass

    for sd in transformed.get("slides",[]):
        render(sd)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─── SESSION STATE ────────────────────────────────────────────────────────────

for key, default in [
    ("logs",           []),
    ("parsed",         None),
    ("design_system",  None),
    ("learned_design", None),
    ("kb_file_names",  []),
    ("transformed",    None),
    ("pptx_bytes",     None),
    ("stage",          "idle"),
    ("system_prompt",  SYSTEM_PROMPT),
]:
    if key not in st.session_state:
        st.session_state[key] = default


# ─── SIDEBAR ─────────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown("""
    <div style="padding:1.5rem 0 1rem;">
      <div style="font-family:'DM Serif Display',serif;font-size:1.3rem;color:#fff;margin-bottom:.25rem;">
        ⚡ PPT Intelligence
      </div>
      <div style="font-size:.72rem;color:#6b7280;font-family:'DM Mono',monospace;">
        Qwen2.5-7B · llama.cpp · n_ctx=8192
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="card-header">LLM ENDPOINT</div>', unsafe_allow_html=True)
    llm_endpoint = st.text_input("endpoint", value="http://localhost:8080/v1/chat/completions",
                                 label_visibility="collapsed")
    llm_model    = st.text_input("model",    value="qwen2.5-coder-7b",
                                 label_visibility="collapsed")

    st.markdown('<div class="card-header" style="margin-top:1rem;">PARAMETERS</div>',
                unsafe_allow_html=True)
    temperature = st.slider("Temperature", 0.0, 1.0, 0.15, 0.05)
    max_tokens  = st.slider("Max output tokens", 500, 3000, 2000, 100,
                            help="Model ctx=8192 total. Input ~1500-2500 tokens per batch. 2000 output is safe.")
    batch_size  = st.slider("Slides per batch", 2, 8, 4, 1,
                            help="4 = default safe. Drop to 2-3 if you still see truncation.")

    st.markdown("---")

    if st.button("🔌 Test Connection"):
        try:
            r = requests.post(llm_endpoint,
                json={"model":llm_model,"messages":[{"role":"user","content":"Reply OK"}],"max_tokens":5},
                timeout=10)
            st.success("✓ Connected") if r.status_code==200 else st.error("HTTP {}".format(r.status_code))
        except Exception as e:
            st.error("✗ {}".format(str(e)[:60]))

    st.markdown("---")

    with st.expander("✏️ Edit System Prompt"):
        edited = st.text_area("prompt", value=st.session_state.system_prompt,
                              height=300, label_visibility="collapsed")
        c1,c2 = st.columns(2)
        if c1.button("💾 Save",  use_container_width=True):
            st.session_state.system_prompt = edited; st.success("Saved")
        if c2.button("↩ Reset", use_container_width=True):
            st.session_state.system_prompt = SYSTEM_PROMPT; st.success("Reset")

    st.markdown("---")

    kb_dot   = "ok"  if st.session_state.learned_design else "off"
    kb_label = "{} ref files · {} slides".format(
        st.session_state.learned_design["files_analysed"],
        st.session_state.learned_design["total_slides_analysed"]
    ) if st.session_state.learned_design else "No KB loaded"

    stage_map = {"idle":("off","Idle"),"parsed":("warn","PPTX parsed"),
                 "transformed":("ok","Transform done"),"done":("ok","PPTX generated")}
    dot, label = stage_map.get(st.session_state.stage, ("off","Idle"))

    st.markdown("""
    <div class="status-bar" style="flex-direction:column;gap:.4rem;">
      <div class="status-item">
        <div class="status-dot {dot}"></div><span>{label}</span>
      </div>
      <div class="status-item">
        <div class="status-dot {kb_dot}"></div>
        <span style="font-size:.68rem;">KB: {kb_label}</span>
      </div>
      <div style="font-size:.68rem;color:#4b5563;margin-top:.2rem;">
        Slides parsed: {sp}<br>Images: {img}<br>Output slides: {out}
      </div>
    </div>
    """.format(
        dot=dot, label=label, kb_dot=kb_dot, kb_label=kb_label,
        sp=len(st.session_state.parsed["slides"]) if st.session_state.parsed else 0,
        img=len(st.session_state.parsed["images"]) if st.session_state.parsed else 0,
        out=len(st.session_state.transformed["slides"]) if st.session_state.transformed else 0,
    ), unsafe_allow_html=True)

    if st.button("🔄 Reset All"):
        for k in ["logs","parsed","design_system","learned_design",
                  "kb_file_names","transformed","pptx_bytes"]:
            st.session_state[k] = [] if k in ("logs","kb_file_names") else None
        st.session_state.stage = "idle"
        st.rerun()


# ─── MAIN UI ─────────────────────────────────────────────────────────────────

st.markdown("""
<div class="hero">
  <div class="hero-badge">AI PRESENTATION OPERATING SYSTEM</div>
  <div class="hero-title">Presentation Intelligence Platform</div>
  <div class="hero-sub">Upload reference decks · Upload target · AI transforms · Download</div>
</div>
""", unsafe_allow_html=True)

# ── 01 KNOWLEDGE BASE ─────────────────────────────────────────────────────────

st.markdown('<div class="section-label">01 · REFERENCE KNOWLEDGE BASE</div>', unsafe_allow_html=True)
col_kb, col_kb_info = st.columns([2,1])

with col_kb:
    kb_files = st.file_uploader("Reference PPTXs", type=["pptx"],
                                accept_multiple_files=True,
                                label_visibility="collapsed", key="kb_uploader")
with col_kb_info:
    st.markdown("""
    <div class="card">
      <div class="card-header">WHY THIS MATTERS</div>
      <ul class="bullet-list">
        <li>Learns your brand colors &amp; fonts</li>
        <li>Infers layout preferences</li>
        <li>More files = stronger signal</li>
        <li>Tip: upload 3–10 best decks</li>
      </ul>
    </div>""", unsafe_allow_html=True)

if kb_files:
    new_names = sorted([f.name for f in kb_files])
    if new_names != st.session_state.kb_file_names:
        with st.spinner("Learning design system from {} file(s)…".format(len(kb_files))):
            try:
                learned = learn_design_system([f.read() for f in kb_files])
                st.session_state.learned_design = learned
                st.session_state.kb_file_names  = new_names
                log("KB learned: {} files, {} slides".format(
                    learned["files_analysed"], learned["total_slides_analysed"]), "ok")
                log("Primary accent: #{}".format(learned["color_palette"]["primary_accent"]), "info")
                log("Fonts: {} / {}".format(
                    learned["typography"]["recommended_pair"]["display"],
                    learned["typography"]["recommended_pair"]["body"]), "info")
                for r in learned["design_rules_inferred"]:
                    log("Rule: {}".format(r), "info")
            except Exception as e:
                log("KB error: {}".format(e), "err")
                st.error("❌ {}".format(e))

if st.session_state.learned_design:
    ld = st.session_state.learned_design
    cp, ty, lp = ld["color_palette"], ld["typography"], ld["layout_patterns"]
    swatches = "".join(
        '<div title="#{c}" style="width:20px;height:20px;border-radius:3px;'
        'background:#{c};border:1px solid rgba(255,255,255,.1);display:inline-block;margin:2px;"></div>'.format(c=c)
        for c in cp["accent_colors"][:8]
    )
    st.markdown("""
    <div class="card" style="border-color:rgba(79,142,247,.3);">
      <div class="card-header">✓ LEARNED — {files} files · {slides} slides</div>
      <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:1rem;">
        <div>
          <div style="font-size:.72rem;color:#6b7280;margin-bottom:.4rem;">PALETTE</div>
          {swatches}
          <div style="font-size:.72rem;color:#9ca3af;margin-top:.4rem;">
            Primary: <code style="color:#4f8ef7;">#{primary}</code><br>
            Secondary: <code style="color:#4f8ef7;">#{secondary}</code>
          </div>
        </div>
        <div>
          <div style="font-size:.72rem;color:#6b7280;margin-bottom:.4rem;">TYPOGRAPHY</div>
          <div style="font-size:.8rem;color:#e8eaf0;">
            <span style="color:#4f8ef7;">Title:</span> {title_font} @ {title_pt}pt<br>
            <span style="color:#4f8ef7;">Body:</span>  {body_font} @ {body_pt}pt
          </div>
        </div>
        <div>
          <div style="font-size:.72rem;color:#6b7280;margin-bottom:.4rem;">LAYOUT MIX</div>
          {layout_rows}
        </div>
      </div>
      <div style="margin-top:.75rem;border-top:1px solid #2a2f42;padding-top:.75rem;">
        <div style="font-size:.7rem;color:#6b7280;margin-bottom:.3rem;">INFERRED RULES</div>
        {rules}
      </div>
    </div>""".format(
        files=ld["files_analysed"], slides=ld["total_slides_analysed"],
        swatches=swatches,
        primary=cp["primary_accent"] or "N/A",
        secondary=cp["secondary_accent"] or "N/A",
        title_font=ty["recommended_pair"]["display"], title_pt=ty["title_size_pt"],
        body_font=ty["recommended_pair"]["body"],   body_pt=ty["body_size_pt"],
        layout_rows="".join(
            '<div style="font-size:.72rem;color:#9ca3af;line-height:1.8;">'
            '{}: <span style="color:#4f8ef7;">{}</span></div>'.format(p,v)
            for p,v in list(lp["distribution"].items())[:5]
        ),
        rules="".join(
            '<div style="font-size:.75rem;color:#c8cdd8;padding:.1rem 0;">▸ {}</div>'.format(r)
            for r in ld["design_rules_inferred"]
        ),
    ), unsafe_allow_html=True)
else:
    st.markdown("""
    <div style="padding:.75rem 1rem;background:rgba(245,158,11,.08);border:1px solid rgba(245,158,11,.2);
         border-radius:8px;font-size:.8rem;color:#fbbf24;margin-top:.5rem;">
      ⚠️ No KB loaded — AI will use default design judgment.
    </div>""", unsafe_allow_html=True)

st.markdown("---")

# ── 02 UPLOAD TARGET ──────────────────────────────────────────────────────────

st.markdown('<div class="section-label">02 · UPLOAD TARGET PRESENTATION</div>', unsafe_allow_html=True)
col_up, col_info = st.columns([2,1])
with col_up:
    uploaded = st.file_uploader("Target PPTX", type=["pptx"], label_visibility="collapsed")
with col_info:
    st.markdown("""
    <div class="card">
      <div class="card-header">WHAT HAPPENS</div>
      <ul class="bullet-list">
        <li>Parse all slides, text &amp; images</li>
        <li>Apply learned design system</li>
        <li>Batch-transform via Qwen2.5</li>
        <li>Generate executive-grade .pptx</li>
      </ul>
    </div>""", unsafe_allow_html=True)

if uploaded and st.session_state.stage == "idle":
    with st.spinner("Parsing…"):
        try:
            fb = uploaded.read()
            st.session_state.parsed       = parse_pptx(fb)
            st.session_state.design_system = extract_design_system(fb)
            st.session_state.stage        = "parsed"
            log("Parsed {} slides, {} images".format(
                st.session_state.parsed["slide_count"],
                len(st.session_state.parsed["images"])), "ok")
        except Exception as e:
            log("Parse error: {}".format(e), "err")
            st.error("❌ Failed to parse PPTX: {}".format(e))

# ── 03 PARSED SUMMARY ─────────────────────────────────────────────────────────

if st.session_state.parsed:
    parsed = st.session_state.parsed
    ds     = st.session_state.design_system
    st.markdown('<div class="section-label">03 · PARSED CONTENT</div>', unsafe_allow_html=True)

    c1,c2,c3,c4 = st.columns(4)
    for col,val,lbl in [
        (c1, parsed["slide_count"],       "Slides"),
        (c2, len(parsed["images"]),        "Images"),
        (c3, ds.get("aspect_ratio","?"),  "Aspect"),
        (c4, len(ds.get("all_fonts",[])), "Fonts"),
    ]:
        col.markdown("""
        <div class="metric-card" style="text-align:center;padding:1rem;">
          <div class="metric-value">{}</div>
          <div class="metric-label">{}</div>
        </div>""".format(val, lbl), unsafe_allow_html=True)

    with st.expander("📋 View parsed slides"):
        for s in parsed["slides"]:
            st.markdown("""
            <div class="slide-card">
              <div class="slide-num">SLIDE {n}</div>
              <div class="slide-title">{title}</div>
              <ul class="bullet-list">{bullets}</ul>
              {imgs}
            </div>""".format(
                n=s["slide_number"],
                title=s["title"] or "(no title)",
                bullets="".join("<li>{}</li>".format(b) for b in s["bullets"][:5]),
                imgs='<div style="margin-top:.5rem;font-size:.72rem;color:#4f8ef7;'
                     'font-family:DM Mono,monospace;">🖼 {}</div>'.format(
                         ", ".join(s["image_ids"])) if s["image_ids"] else "",
            ), unsafe_allow_html=True)

# ── 04 AI TRANSFORM ───────────────────────────────────────────────────────────

if st.session_state.stage == "parsed":
    st.markdown('<div class="section-label">04 · AI TRANSFORM</div>', unsafe_allow_html=True)

    if st.session_state.learned_design:
        ld = st.session_state.learned_design
        st.markdown("""
        <div style="padding:.6rem 1rem;background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);
             border-radius:8px;font-size:.8rem;color:#4ade80;margin-bottom:.75rem;">
          ✓ Design KB active — {files} ref files · Primary #{primary} ·
          Fonts: {tf} / {bf}
        </div>""".format(
            files=ld["files_analysed"],
            primary=ld["color_palette"]["primary_accent"],
            tf=ld["typography"]["recommended_pair"]["display"],
            bf=ld["typography"]["recommended_pair"]["body"],
        ), unsafe_allow_html=True)
    else:
        st.markdown("""
        <div style="padding:.6rem 1rem;background:rgba(245,158,11,.08);border:1px solid rgba(245,158,11,.2);
             border-radius:8px;font-size:.8rem;color:#fbbf24;margin-bottom:.75rem;">
          ⚠️ No KB — AI will use default design judgment. Upload reference PPTXs in Step 01.
        </div>""", unsafe_allow_html=True)

    if st.button("⚡ Transform with Qwen2.5", use_container_width=True):
        n_slides  = st.session_state.parsed["slide_count"]
        n_batches = max(1, -(-n_slides // batch_size))
        log("Starting: {} slides → {} batch(es) of ≤{}".format(n_slides, n_batches, batch_size), "info")

        progress   = st.progress(0, text="Starting…")
        status_box = st.empty()

        def upd(pct, msg):
            progress.progress(pct, text=msg)
            status_box.markdown(
                '<div style="font-size:.78rem;color:#6b7280;font-family:DM Mono,monospace;">{}</div>'.format(msg),
                unsafe_allow_html=True)

        try:
            result = transform_in_batches(
                parsed         = st.session_state.parsed,
                file_design    = st.session_state.design_system,
                learned_design = st.session_state.learned_design,
                endpoint       = llm_endpoint,
                model          = llm_model,
                temperature    = temperature,
                max_tokens     = max_tokens,
                system_prompt  = st.session_state.system_prompt,
                batch_size     = batch_size,
                progress_cb    = upd,
            )
            upd(100, "Done!")
            status_box.empty()

            if result.get("_truncated"):
                st.warning("⚠️ Some batches were truncated and auto-repaired. "
                           "Reduce 'Slides per batch' to 2–3 for cleaner results.")

            st.session_state.transformed = result
            st.session_state.stage = "transformed"
            log("Complete — {} output slides".format(len(result.get("slides",[]))), "ok")
            st.rerun()

        except requests.exceptions.ConnectionError:
            st.error("❌ Cannot connect — is llama.cpp running?")
            progress.empty(); status_box.empty()
        except json.JSONDecodeError as e:
            st.error("❌ JSON failed after repair. Reduce batch size to 2.\n\n`{}`".format(str(e)[:200]))
            log(str(e), "err")
            progress.empty(); status_box.empty()
        except Exception as e:
            st.error("❌ {}".format(e))
            log(str(e), "err")
            progress.empty(); status_box.empty()

# ── 05 RESULTS ────────────────────────────────────────────────────────────────

if st.session_state.transformed:
    t = st.session_state.transformed
    st.markdown('<div class="section-label">05 · TRANSFORMED DECK</div>', unsafe_allow_html=True)
    st.markdown("""
    <div class="deck-header">
      <div class="deck-title">{title}</div>
      <div class="deck-meta">{out} slides · Theme: {theme} · Original: {orig} slides</div>
    </div>""".format(
        title=t.get("deck_title","Transformed Deck"),
        out=len(t.get("slides",[])),
        theme=t.get("deck_theme","dark"),
        orig=st.session_state.parsed["slide_count"],
    ), unsafe_allow_html=True)

    for slide in t.get("slides",[]):
        layout     = slide.get("layout","bullet")
        bullets    = slide.get("bullets",[])
        flow_steps = slide.get("flow_steps",[])
        metrics    = slide.get("metrics",[])
        lc         = slide.get("left_column",{})
        rc         = slide.get("right_column",{})

        flow_html = ""
        if flow_steps:
            items = ["".join(
                '<span class="flow-arrow">→</span>' if s in ("->","→")
                else '<span class="flow-step">{}</span>'.format(s)
                for s in flow_steps)]
            flow_html = '<div class="flow-steps">{}</div>'.format("".join(items))

        metrics_html = ""
        if metrics:
            cards = "".join(
                '<div class="metric-card">'
                '<div class="metric-value">{}</div>'
                '<div class="metric-label">{}</div></div>'.format(m.get("value",""),m.get("label",""))
                for m in metrics)
            metrics_html = '<div class="metric-grid">{}</div>'.format(cards)

        two_col_html = ""
        if layout == "two-column":
            two_col_html = """
            <div class="two-col">
              <div class="col-box">
                <div class="col-heading">{lh}</div>
                <ul class="bullet-list">{lb}</ul>
              </div>
              <div class="col-box">
                <div class="col-heading">{rh}</div>
                <ul class="bullet-list">{rb}</ul>
              </div>
            </div>""".format(
                lh=lc.get("heading",""), lb="".join("<li>{}</li>".format(b) for b in lc.get("bullets",[])),
                rh=rc.get("heading",""), rb="".join("<li>{}</li>".format(b) for b in rc.get("bullets",[])),
            )

        st.markdown("""
        <div class="slide-card">
          <div class="slide-num">SLIDE {id} · {layout}</div>
          <div class="slide-title">{title}</div>
          {sub}
          <div class="slide-meta">
            <span class="tag tag-class">{cls}</span>
            <span class="tag tag-layout">{layout}</span>
            <span class="tag tag-visual">{vis}</span>
            {img_tag}
          </div>
          {bullets_html}
          {flow_html}
          {metrics_html}
          {two_col_html}
          {notes_html}
        </div>""".format(
            id=slide.get("id",""),
            layout=layout,
            title=slide.get("title",""),
            sub='<div style="font-size:.82rem;color:#9ca3af;margin-bottom:.5rem;">{}</div>'.format(
                slide.get("subtitle","")) if slide.get("subtitle") else "",
            cls=slide.get("classification",""),
            vis=slide.get("visual_type",""),
            img_tag='<span class="tag tag-visual">🖼 {} img</span>'.format(
                len(slide.get("images",[]))) if slide.get("images") else "",
            bullets_html="<ul class='bullet-list'>{}</ul>".format(
                "".join("<li>{}</li>".format(b) for b in bullets)) if bullets else "",
            flow_html=flow_html,
            metrics_html=metrics_html,
            two_col_html=two_col_html,
            notes_html='<div style="margin-top:.75rem;font-size:.75rem;color:#6b7280;'
                       'font-style:italic;border-left:2px solid #2a2f42;padding-left:.75rem;">'
                       '{}</div>'.format(slide["notes"]) if slide.get("notes") else "",
        ), unsafe_allow_html=True)

    with st.expander("🔍 Raw JSON"):
        st.json(t)

# ── 06 GENERATE ───────────────────────────────────────────────────────────────

if st.session_state.transformed and st.session_state.stage == "transformed":
    st.markdown('<div class="section-label">06 · GENERATE & DOWNLOAD</div>', unsafe_allow_html=True)
    col_gen, col_dl = st.columns(2)

    with col_gen:
        if st.button("🏗️ Generate .pptx", use_container_width=True):
            with st.spinner("Building PowerPoint…"):
                try:
                    pptx_bytes = generate_pptx(
                        st.session_state.transformed,
                        st.session_state.parsed["images"],
                    )
                    st.session_state.pptx_bytes = pptx_bytes
                    st.session_state.stage = "done"
                    log("PPTX generated — {:,} bytes".format(len(pptx_bytes)), "ok")
                    st.rerun()
                except Exception as e:
                    st.error("❌ {}\n\n{}".format(e, traceback.format_exc()))

    with col_dl:
        if st.session_state.pptx_bytes:
            deck_title = st.session_state.transformed.get("deck_title","deck")
            safe_name  = "".join(c if c.isalnum() or c in " -_" else "_" for c in deck_title)
            st.download_button(
                "⬇️ Download Executive Deck",
                data=st.session_state.pptx_bytes,
                file_name="{}.pptx".format(safe_name),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

# ── LOG ───────────────────────────────────────────────────────────────────────

if st.session_state.logs:
    st.markdown('<div class="section-label">PIPELINE LOG</div>', unsafe_allow_html=True)
    render_logs()
